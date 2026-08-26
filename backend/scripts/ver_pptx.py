"""Dibuja una hoja de un PPTX como PNG, para poder MIRARLA.

No hay PowerPoint ni LibreOffice en esta maquina, asi que no se puede convertir
de verdad. Esto reconstruye la hoja con PIL a partir de lo que el archivo
declara: la imagen de fondo del master, y cada cuadro de texto en su posicion,
con su tamano, su color y su alineacion.

Es una APROXIMACION y hay que tratarla como tal: PowerPoint parte las lineas y
ajusta el interletrado distinto. Sirve para lo que importa acá --ver si algo se
sale, si dos cosas se pisan, si un tamano quedo desparejo-- y no para juzgar
como queda un remate fino.

La fuente real ayuda: el diseno usa Calibri, que si esta en esta maquina, asi
que el ancho medido es el de verdad y no una estimacion.
"""
import io, os, sys
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

CM = 360000
FUENTES = {
    ("calibri", False): r"C:\Windows\Fonts\calibri.ttf",
    ("calibri", True):  r"C:\Windows\Fonts\calibrib.ttf",
    ("arial", False):   r"C:\Windows\Fonts\arial.ttf",
    ("arial", True):    r"C:\Windows\Fonts\arialbd.ttf",
}
# Las que no estan instaladas se dibujan con Calibri y se avisa al final.
_FALTANTES: set[str] = set()


def _fuente(nombre, bold, px):
    clave = (str(nombre or "calibri").strip().lower(), bool(bold))
    ruta = FUENTES.get(clave)
    if ruta is None:
        _FALTANTES.add(str(nombre))
        ruta = FUENTES[("calibri", bool(bold))]
    return ImageFont.truetype(ruta, max(1, int(px)))


def _partir(texto, fuente, ancho_px):
    """Igual que hace PowerPoint: parte por palabra al llegar al borde."""
    lineas, actual = [], ""
    for palabra in texto.split():
        prueba = f"{actual} {palabra}".strip()
        if fuente.getlength(prueba) <= ancho_px or not actual:
            actual = prueba
        else:
            lineas.append(actual); actual = palabra
    if actual: lineas.append(actual)
    return lineas or [""]


def dibujar(ruta_pptx, hoja=0, salida=None, ppcm=38):
    prs = Presentation(ruta_pptx)
    W = int(prs.slide_width / CM * ppcm)
    H = int(prs.slide_height / CM * ppcm)
    lienzo = Image.new("RGB", (W, H), "white")

    # El fondo: la imagen del master, que es donde vive el arte de la cenefa.
    try:
        arte = prs.slides[hoja].slide_layout.slide_master.shapes[0]
        fondo = Image.open(io.BytesIO(arte.image.blob)).convert("RGB").resize((W, H))
        lienzo.paste(fondo, (0, 0))
    except Exception as e:
        print(f"   (sin arte de fondo: {e})")

    d = ImageDraw.Draw(lienzo)

    # Primero los rellenos. Sin esto la cocarda no se ve: su texto es BLANCO y
    # sin el rectangulo rojo abajo queda blanco sobre blanco, o sea invisible
    # --y uno concluye que la cocarda no se dibujo, que es exactamente el tipo
    # de conclusion equivocada que este visor tiene que evitar.
    for s in prs.slides[hoja].shapes:
        try:
            if s.fill.type is None or s.fill.type == 5:   # 5 = hereda del fondo
                continue
            rgb = s.fill.fore_color.rgb
        except Exception:
            continue
        caja = [(s.left or 0)/CM*ppcm, (s.top or 0)/CM*ppcm,
                ((s.left or 0)+(s.width or 0))/CM*ppcm, ((s.top or 0)+(s.height or 0))/CM*ppcm]
        d.rounded_rectangle(caja, radius=ppcm*0.25, fill="#%02X%02X%02X" % tuple(rgb))

    # Las lineas (el tachado del precio anterior). Sin esto no se ve si la regla
    # de "con mecanica no se tacha" hizo algo o no.
    for s in prs.slides[hoja].shapes:
        if "LINE" not in str(s.shape_type):
            continue
        x0, y0 = (s.left or 0)/CM*ppcm, (s.top or 0)/CM*ppcm
        x1, y1 = x0 + (s.width or 0)/CM*ppcm, y0 + (s.height or 0)/CM*ppcm
        # Las de PowerPoint van de abajo-izquierda a arriba-derecha.
        d.line([(x0, y1), (x1, y0)], fill="#111111", width=max(2, int(ppcm*0.06)))

    for s in sorted(prs.slides[hoja].shapes, key=lambda s: (s.top or 0)):
        if not (s.has_text_frame and s.text_frame.text.strip()):
            continue
        x0 = (s.left or 0) / CM * ppcm
        y = (s.top or 0) / CM * ppcm
        ancho = (s.width or 0) / CM * ppcm
        alto  = (s.height or 0) / CM * ppcm

        # Anclaje vertical. PowerPoint centra o abaja el bloque de texto dentro
        # de la caja; dibujarlo siempre desde arriba corre todo y hace ver
        # solapamientos que no existen (y tapa los que si).
        bloque = 0
        for p in s.text_frame.paragraphs:
            rr = [r for r in p.runs if r.text]
            if rr:
                pt0 = rr[0].font.size.pt if rr[0].font.size else 18
                f0 = _fuente(rr[0].font.name, rr[0].font.bold, pt0/72*2.54*ppcm)
                bloque += len(_partir("".join(r.text for r in rr), f0, ancho or W)) * pt0/72*2.54*ppcm*1.2
        anc = str(s.text_frame.vertical_anchor or "").lower()
        if "middle" in anc:   y += max(0, (alto - bloque) / 2)
        elif "bottom" in anc: y += max(0, alto - bloque)
        for p in s.text_frame.paragraphs:
            runs = [r for r in p.runs if r.text]
            if not runs: continue
            alin = str(p.alignment or "").lower()

            # Run por run, no todo el parrafo con el tamano del primero: el
            # cuadro del precio son TRES runs de tamanos muy distintos ("$" a
            # 71 pt, "149" a 128, los centavos a 28). Medirlo todo con 71
            # dibuja el precio a la mitad de lo que sale impreso.
            piezas = []
            for r in runs:
                pt = r.font.size.pt if r.font.size else 18
                px = pt / 72 * 2.54 * ppcm
                f = _fuente(r.font.name, r.font.bold, px)
                try:
                    col = "#%02X%02X%02X" % tuple(r.font.color.rgb) if r.font.color and r.font.color.rgb else "#000000"
                except Exception:
                    col = "#000000"
                piezas.append((r.text, f, px, col))

            # Una sola linea si entra; si no, se parte por el run mas ancho.
            total = sum(pz[1].getlength(pz[0]) for pz in piezas)
            alto_linea = max(pz[2] for pz in piezas) * 1.2
            if total <= (ancho or W) or len(piezas) > 1:
                if "center" in alin:   lx = x0 + (ancho - total) / 2
                elif "right" in alin:  lx = x0 + ancho - total
                else:                  lx = x0
                base = y + alto_linea
                for texto, f, px, col in piezas:
                    # Los runs chicos se alinean por ARRIBA de los grandes: asi
                    # es como se dibujan los centavos del precio.
                    d.text((lx, base - px * 1.2), texto, font=f, fill=col)
                    lx += f.getlength(texto)
                y += alto_linea
            else:
                texto, f, px, col = piezas[0][0], piezas[0][1], piezas[0][2], piezas[0][3]
                for linea in _partir(texto, f, ancho or W):
                    lw = f.getlength(linea)
                    if "center" in alin:   lx = x0 + (ancho - lw) / 2
                    elif "right" in alin:  lx = x0 + ancho - lw
                    else:                  lx = x0
                    d.text((lx, y), linea, font=f, fill=col)
                    y += px * 1.2
    salida = salida or os.path.splitext(ruta_pptx)[0] + f"_hoja{hoja+1}.png"
    lienzo.save(salida)
    if _FALTANTES:
        print(f"   (dibujadas con Calibri por no estar instaladas: {sorted(_FALTANTES)})")
    return salida


if __name__ == "__main__":
    print(dibujar(sys.argv[1], int(sys.argv[2]) - 1 if len(sys.argv) > 2 else 0,
                  sys.argv[3] if len(sys.argv) > 3 else None))
