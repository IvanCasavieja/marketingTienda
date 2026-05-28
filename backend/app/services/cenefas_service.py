"""Cenefas PPTX generator — ported from generar_cenefas.py."""
import copy
import io
import re
import unicodedata
from typing import Any

import openpyxl
from lxml import etree
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

# ---------------------------------------------------------------------------
# Format constants
# ---------------------------------------------------------------------------

P1_FONT_SIZE  = 32      # pt — "Precio Final" / "6X" label
P1_BOLD       = False   # label goes without bold
P1_MARGIN_EMU = 466400  # distance from P1 top to price shape top (32pt*12700 + 60000 gap)

PRICE_DECIMAL_PT  = 20  # pt — parte decimal fija (,90 / ,20) para todos los precios
PBANCO_INT_PT     = 40  # pt — tamaño fijo del entero/símbolo de pBanco
UNIDAD_PRECIO_PT  = 16  # pt — "unidad" debajo de Precio1 en productos multi-SKU
UNIDAD_PBANCO_PT  = 10  # pt — "unidad" debajo de pBanco en productos multi-SKU

DELI_SUBCATS = {"FIAMBRES", "QUESOS"}
NO_UNIDAD_SUBCATS = {"CARNES", "FIAMBRES", "EMBUTIDOS CARNE", "QUESOS"}


# ---------------------------------------------------------------------------
# Price formatting
# ---------------------------------------------------------------------------

def fmt_price(value: Any) -> str:
    if value is None:
        return "0"
    value = float(value)
    if value == int(value):
        val = int(value)
        return f"{val:,}".replace(",", ".") if val >= 1000 else str(val)
    int_part = int(value)
    dec_str = f"{value:.2f}".split(".")[1]
    int_str = f"{int_part:,}".replace(",", ".") if int_part >= 1000 else str(int_part)
    return int_str + "," + dec_str


def _parse_price_raw(value: Any) -> float:
    """Convierte un valor de precio crudo a float.

    Soporta: números (int/float), strings con texto como '181 UNIDAD',
    decimales con coma ('144,8') y separadores de miles con punto ('1.500').
    """
    if value is None or value == "":
        return 0.0
    if isinstance(value, (int, float)):
        return float(value)
    m = re.match(r"^(\d[\d.,]*)", str(value).strip())
    if not m:
        return 0.0
    num_str = m.group(1)
    # "1.500,50" → European: punto=miles, coma=decimal
    if "." in num_str and "," in num_str:
        num_str = num_str.replace(".", "").replace(",", ".")
    elif "," in num_str:
        num_str = num_str.replace(",", ".")
    try:
        return float(num_str)
    except ValueError:
        return 0.0


# ---------------------------------------------------------------------------
# Combo parsing
# ---------------------------------------------------------------------------

def parse_combo(oferta_str: str) -> tuple[str, float]:
    s = str(oferta_str or "").strip()
    m = re.match(r"(\d+)\s*[xX]\s*\$?\s*(\d+(?:[.,]\d+)?)\s*$", s)
    if m:
        amount = float(m.group(2).replace(",", "."))
        return m.group(1) + "X", amount
    return "", 0.0


# ---------------------------------------------------------------------------
# Bold detection for brand names in ALL-CAPS
# ---------------------------------------------------------------------------

def split_caps(text: str) -> list[tuple[str, bool]]:
    pattern = r"([A-Z]{2,}[A-Z0-9\-/]*(?:\s+[A-Z]{2,}[A-Z0-9\-/]*)*)"
    parts = re.split(pattern, text)
    result = []
    for part in parts:
        if not part:
            continue
        is_bold = bool(re.fullmatch(
            r"[A-Z]{2,}[A-Z0-9\-/]*(?:\s+[A-Z]{2,}[A-Z0-9\-/]*)*", part
        ))
        result.append((part, is_bold))
    return result


# ---------------------------------------------------------------------------
# Row processing
# ---------------------------------------------------------------------------

def process_row(row: tuple, h: dict, vigencia: str, aclaracion: str, otra_alcohol: str, banco: str = "") -> dict:
    ofertadet = str(row[h["OFERTADET"]] or "").strip() if "OFERTADET" in h else "Precio fijo"
    precio_raw = row[h["PRECIO"]] if "PRECIO" in h else 0
    oferta_raw = row[h["OFERTA"]] if "OFERTA" in h else ""
    desc = str(row[h["DESCRIPCION"]] or "").strip()
    cat = str(row[h["Categoria"]] or "").strip() if "Categoria" in h else ""
    subcat = str(row[h["subcategoria"]] or "").strip() if "subcategoria" in h else ""
    moneda = str(row[h["MONEDA"]] or "").strip() if "MONEDA" in h else "$"
    code = str(row[h["CODIGO"]] or "").strip() if "CODIGO" in h else ""
    precio_banco_raw = row[h["PRECIO_BANCO"]] if "PRECIO_BANCO" in h else None

    prefix = "U$S " if moneda == "U$S" else "$"
    precio = _parse_price_raw(precio_raw)

    p1 = ""
    precio_display = ""
    mecanica = ""

    if ofertadet == "Combo":
        p1, amount = parse_combo(oferta_raw)
        precio_display = prefix + fmt_price(amount)
        qty = p1[:-1] if p1.endswith("X") else "2"
        mecanica = f"Comprando {qty}, {prefix}{fmt_price(precio)} c/u"

    elif ofertadet == "M x N":
        precio_display = prefix + fmt_price(precio)
        mecanica = f"Comprando 2, {prefix}{fmt_price(precio)} c/u"

    else:
        precio_val = precio
        if subcat in DELI_SUBCATS:
            dl = desc.lower()
            if ". kg" in dl or " kg" in dl or dl.endswith("kg") or "100g" in dl:
                precio_val = precio / 10
                desc = re.sub(r"\.\s*[Kk]g\b", ". 100g", desc)
                desc = re.sub(r"\s+[Kk]g\b", " 100g", desc)
        precio_display = prefix + fmt_price(precio_val)

    if not p1:
        p1 = "Precio Final"

    # "unidad" solo para productos con múltiples SKUs (código con "/" o "dígito - dígito").
    # Productos de un solo SKU no muestran "unidad" aunque sean precio fijo.
    is_multi_sku = bool(code and ("/" in code or re.search(r"\d\s*[-–—]\s*\d", code)))

    if is_multi_sku and ofertadet in ("Precio fijo", "% descuento"):
        unidad = "" if subcat in NO_UNIDAD_SUBCATS else "unidad"
    else:
        unidad = ""

    otra = otra_alcohol if cat == "BEBIDAS CON ALCOHOL" else ""

    # Precio bancario: formateado igual que precio principal
    pbanco_display = ""
    if precio_banco_raw is not None and precio_banco_raw != "":
        pbanco_val = _parse_price_raw(precio_banco_raw)
        if pbanco_val > 0:
            pbanco_display = prefix + fmt_price(pbanco_val)

    unidad_precio = unidad
    unidad_pbanco = "unidad" if is_multi_sku else ""

    return {
        "p1": p1,
        "precio": precio_display,
        "mecanica": mecanica,
        "descripcion": desc,
        "unidad": unidad,
        "vigencia": vigencia,
        "aclaracion": aclaracion,
        "otra_aclaracion": otra,
        "code": code,
        "pbanco": pbanco_display,
        "banco": banco,
        "unidad_precio": unidad_precio,
        "unidad_pbanco": unidad_pbanco,
    }


def _normalize_header(name: str) -> str:
    """Normaliza a mayúsculas sin tildes para matching flexible de columnas."""
    return unicodedata.normalize("NFD", str(name)).encode("ascii", "ignore").decode().upper().strip()


_CANONICAL_COLUMNS = [
    "Categoria", "subcategoria", "OFERTADET", "DESCRIPCION", "PRECIO",
    "OFERTA", "MONEDA", "NOMBREARTICULO", "CODIGO", "PRECIO_BANCO",
]

# Aliases: nombre normalizado del Excel → nombre canónico interno
# Permite columnas con nombres alternativos sin romper el sistema.
_HEADER_ALIASES: dict[str, str] = {
    "PRECIOS": "PRECIO",           # variante plural
    "SCOTLAND 20%": "PRECIO_BANCO",
    "SCOTIA 20%": "PRECIO_BANCO",  # columna de precio bancario Scotia
    "PRECIO BANCO": "PRECIO_BANCO",
    "PBANCO": "PRECIO_BANCO",
}

_EXPECTED_HEADERS: dict[str, str] = {
    _normalize_header(k): k for k in _CANONICAL_COLUMNS
}
_EXPECTED_HEADERS.update({
    _normalize_header(alias): canonical
    for alias, canonical in _HEADER_ALIASES.items()
})

_OPTIONAL_HEADERS = {"Categoria", "subcategoria", "OFERTADET", "OFERTA", "CODIGO", "PRECIO_BANCO", "NOMBREARTICULO"}

# Columnas mínimas que deben existir para que el parsing funcione
_REQUIRED_HEADERS = {"DESCRIPCION", "PRECIO"}

# Columnas usadas para auto-detectar la fila de headers
_DETECTION_COLS = {"OFERTADET", "DESCRIPCION", "PRECIOS", "PRECIO", "CODIGO", "MONEDA"}


def load_products_from_bytes(
    excel_bytes: bytes,
    vigencia: str,
    aclaracion: str,
    otra_alcohol: str,
    banco: str = "",
) -> list[dict]:
    wb = openpyxl.load_workbook(io.BytesIO(excel_bytes))
    ws = wb["Cenefas"] if "Cenefas" in wb.sheetnames else wb.active

    # Auto-detectar fila de headers: buscar la primera fila que contenga
    # al menos 2 columnas conocidas (no solo OFERTADET, para soportar
    # Excels alternativos como el de Especial Mascotas).
    header_row = None
    for i, row in enumerate(ws.iter_rows(max_row=10, values_only=True), start=1):
        normalized = {_normalize_header(str(c)) for c in row if c is not None}
        if "OFERTADET" in normalized or len(normalized & _DETECTION_COLS) >= 2:
            header_row = i
            break
    if header_row is None:
        header_row = 1

    raw_headers = [cell.value for cell in ws[header_row]]
    # Mapeo flexible: normaliza el header del Excel al nombre canónico esperado
    h = {}
    for idx, raw in enumerate(raw_headers):
        if not raw:
            continue
        canonical = _EXPECTED_HEADERS.get(_normalize_header(str(raw)))
        h[canonical or str(raw)] = idx

    for k in _REQUIRED_HEADERS:
        if k not in h:
            raise KeyError(k)

    products = []
    seen: set = set()

    for row in ws.iter_rows(min_row=header_row + 1, values_only=True):
        # Saltar filas vacías usando DESCRIPCION como ancla
        if not row[h["DESCRIPCION"]]:
            continue
        # Saltar también si OFERTADET está y está vacío (comportamiento original)
        if "OFERTADET" in h and not row[h["OFERTADET"]]:
            continue
        data = process_row(row, h, vigencia, aclaracion, otra_alcohol, banco)
        key = (data["p1"], data["precio"], data["mecanica"], data["descripcion"].lower().strip(), data.get("code", ""))
        if key not in seen:
            seen.add(key)
            products.append(data)

    return products


# ---------------------------------------------------------------------------
# PPTX manipulation
# ---------------------------------------------------------------------------

def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)


def _set_text(shape, text: str) -> None:
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        runs = para.runs
        if runs:
            runs[0].text = text
            for run in runs[1:]:   # Limpiar runs adicionales (placeholder multi-run)
                run.text = ""
            return
    if shape.text_frame.paragraphs:
        para = shape.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = text


def _set_text_sized(shape, text: str, pt: int) -> None:
    """Sets shape text and forces all runs to the given font size."""
    _set_text(shape, text)
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(pt)


def _set_price(shape, text: str, int_pt: int | None = None) -> None:
    """Renderiza un precio con símbolo, entero y decimal en runs separados.

    Símbolo e entero usan el tamaño que tiene el run en el template.
    Solo el decimal se fuerza a PRICE_DECIMAL_PT (20 pt).
    int_pt se mantiene por compatibilidad pero ya no sobreescribe tamaños.
    """
    if not shape.has_text_frame:
        return

    if text.startswith("U$S "):
        symbol, number = "U$S ", text[4:]
    elif text.startswith("$"):
        symbol, number = "$", text[1:]
    else:
        _set_text(shape, text)
        return

    target_para = None
    for para in shape.text_frame.paragraphs:
        if para.runs:
            target_para = para
            break
    if target_para is None:
        return

    p_elem = target_para._p
    tmpl_r = copy.deepcopy(target_para.runs[0]._r)

    for r_elem in list(p_elem.findall(qn("a:r"))):
        p_elem.remove(r_elem)

    # Split number into integer and decimal parts (decimal separator is ',')
    if "," in number:
        num_int, num_dec_digits = number.rsplit(",", 1)
        num_dec = "," + num_dec_digits
    else:
        num_int = number
        num_dec = None

    # Símbolo e entero: respetan el tamaño escrito en el template.
    # Decimal: siempre fijo en PRICE_DECIMAL_PT (20 pt).
    sym_r = copy.deepcopy(tmpl_r)
    sym_r.find(qn("a:t")).text = symbol

    int_r = copy.deepcopy(tmpl_r)
    int_r.find(qn("a:t")).text = num_int

    runs = [sym_r, int_r]

    if num_dec:
        dec_r = copy.deepcopy(tmpl_r)
        dec_r.find(qn("a:t")).text = num_dec
        dec_rPr = dec_r.find(qn("a:rPr"))
        if dec_rPr is not None:
            dec_rPr.set("sz", str(PRICE_DECIMAL_PT * 100))
        runs.append(dec_r)

    # Insert runs BEFORE <a:endParaRPr> to maintain valid OOXML element order.
    # Appending after endParaRPr causes PowerPoint to silently ignore the runs.
    end_rpr = p_elem.find(qn("a:endParaRPr"))
    if end_rpr is not None:
        idx = list(p_elem).index(end_rpr)
        for r in reversed(runs):
            p_elem.insert(idx, r)
    else:
        for r in runs:
            p_elem.append(r)


def _set_desc(shape, text: str) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    target_para = None
    for para in tf.paragraphs:
        if para.runs:
            target_para = para
            break
    if target_para is None:
        return

    parts = split_caps(text)
    first_seg, first_bold = parts[0]
    first_run = target_para.runs[0]
    tmpl_r = copy.deepcopy(first_run._r)
    first_run.text = first_seg
    first_run.font.bold = True if first_bold else None

    p_elem = target_para._p
    all_r = list(p_elem.findall(qn("a:r")))
    for r_elem in all_r[1:]:
        p_elem.remove(r_elem)

    for seg, bold in parts[1:]:
        if not seg:
            continue
        new_r = copy.deepcopy(tmpl_r)
        new_r.find(qn("a:t")).text = seg
        rPr = new_r.find(qn("a:rPr"))
        if rPr is not None:
            if bold:
                rPr.set("b", "1")
            else:
                rPr.attrib.pop("b", None)
        end_rpr = p_elem.find(qn("a:endParaRPr"))
        if end_rpr is not None:
            p_elem.insert(list(p_elem).index(end_rpr), new_r)
        else:
            p_elem.append(new_r)


def _set_p1(shape, text: str) -> None:
    _set_text(shape, text)
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(P1_FONT_SIZE)
            run.font.bold = P1_BOLD


def _is_multi_sku(code: str) -> bool:
    return bool(code and ("/" in code or re.search(r"\d\s*[-–—]\s*\d", code)))


def _fill_slot(shapes, data: dict, adjust_p1: bool = True) -> None:
    # Expand group shapes so children (e.g. <<banco>>/<<pbanco>>) are reached
    expanded = list(shapes)
    for shape in list(shapes):
        if hasattr(shape, 'shapes'):
            expanded.extend(shape.shapes)

    p1_shape = None
    price_shape = None
    code = data.get("code", "")
    multi = _is_multi_sku(code)

    for shape in expanded:
        t = _shape_text(shape)
        if re.search(r"<<P\d+>>", t):
            p1_shape = shape
            _set_p1(shape, data["p1"])
        elif re.search(r"Precio\s+\d+", t) or re.search(r"<<Precio\d*>>", t):
            price_shape = shape
            _set_price(shape, data["precio"])
        elif re.search(r"<<Mecanica\d+>>", t):
            _set_text(shape, data["mecanica"])
        elif "<<" in t and "Descripci" in t:
            _set_desc(shape, data["descripcion"])
        elif re.search(r"<<UnidadMedida\d+>>", t):
            _set_text(shape, data["unidad"] if multi else "")
        elif re.search(r"<<Vigencia\d*>>", t):
            _set_text(shape, data["vigencia"])
        elif re.search(r"<<Aclaracion\d*>>", t):
            _set_text(shape, data["aclaracion"])
        elif re.search(r"<<OtraAclaracion\d*>>", t):
            _set_text(shape, data["otra_aclaracion"])
        elif re.search(r"<<[Cc]ode\d*>>", t):
            _set_text(shape, code)
        elif re.search(r"<<[Pp][Bb]anco\d*>>", t):
            _set_price(shape, data.get("pbanco", ""), int_pt=PBANCO_INT_PT)
        elif re.search(r"<<[Bb]anco\d*>>", t):
            _set_text(shape, data.get("banco", ""))
        elif re.search(r"<<UnidadPrecio\d*>>", t):
            _set_text_sized(shape, "unidad" if multi else "", UNIDAD_PRECIO_PT)
        elif re.search(r"<<UnidadPBanco\d*>>", t):
            _set_text_sized(shape, "unidad" if multi else "", UNIDAD_PBANCO_PT)
        elif t.strip().lower() == "unidad":
            # Texto estático "unidad" en el template — se limpia para SKU único
            _set_text(shape, "unidad" if multi else "")

    # Ajuste dinámico de P1 solo para plantillas de 1 producto por hoja (A4).
    # En multi-producto, P1 queda en la posición fija de la plantilla.
    if adjust_p1 and p1_shape is not None and price_shape is not None:
        p1_shape.top = price_shape.top - P1_MARGIN_EMU


def _clear_slot(shapes) -> None:
    for shape in shapes:
        if not shape.has_text_frame:
            continue
        t = _shape_text(shape)
        if "<<" in t or re.search(r"Precio\s+\d", t):
            _set_text(shape, "")


def _slot_num(shape) -> int | None:
    """Return the slot number (1–9) if the shape belongs to a numbered product slot."""
    t = _shape_text(shape)
    for n in range(1, 10):
        if (f"<<P{n}>>" in t
                or f"Precio {n}" in t
                or f"<<Precio{n}>>" in t
                or f"<<Mecanica{n}>>" in t
                or f"<<UnidadMedida{n}>>" in t
                or f"<<Vigencia{n}>>" in t
                or f"<<Aclaracion{n}>>" in t
                or f"<<OtraAclaracion{n}>>" in t
                or f"<<UnidadPrecio{n}>>" in t
                or f"<<UnidadPBanco{n}>>" in t
                or ("<<" in t and "Descripci" in t and str(n) in t)):
            return n
    return None


def _get_slots(shapes) -> list[list]:
    """Group shapes by product slot. Supports 1–9 products per slide.

    Strategy:
    - If multiple distinct slot numbers are detected (e.g. <<P1>>/<<P2>>/…), use
      those buckets directly (standard numbered templates, Pinchos 1-9).
    - If all shapes map to slot 1 (templates where every placeholder is labelled
      "1" regardless of position, like Bases cenefas BLACK), find the P1/Price
      anchor shapes, then assign every shape to its spatially nearest anchor.
      This is robust against XML insertion order (shapes added later don't land
      in wrong slots just because they appear at the end of the XML).
    """
    buckets: dict[int, list] = {n: [] for n in range(1, 10)}
    for shape in shapes:
        n = _slot_num(shape)
        if n:
            buckets[n].append(shape)

    present = sorted(n for n in range(1, 10) if buckets[n])

    if len(present) > 1:
        return [buckets[n] for n in present]

    # All shapes land in slot 1 — count anchors to determine product count.
    all_shapes = list(shapes)
    anchors = [s for s in all_shapes if re.search(r"<<P\d+>>", _shape_text(s))]
    if not anchors:
        anchors = [s for s in all_shapes if re.search(r"Precio\s+\d+|<<Precio\d*>>", _shape_text(s))]

    products = len(anchors)
    if products <= 1:
        return [all_shapes]

    # Determine layout orientation from anchor spread (horizontal vs vertical).
    xs = [s.left for s in anchors]
    ys = [s.top for s in anchors]
    horizontal = (max(xs) - min(xs)) >= (max(ys) - min(ys))

    # Sort ALL shapes by spatial center, then split consecutively.
    # Nearest-anchor fails when a shape sits at the far end of its section
    # (e.g. "unidad" below a price) and the NEXT product's anchor is closer.
    # A positional sort + consecutive split is correct for any non-overlapping
    # grid layout regardless of XML insertion order.
    key = (lambda s: s.left + s.width // 2) if horizontal else (lambda s: s.top + s.height // 2)
    sorted_shapes = sorted(all_shapes, key=key)
    per_slot = len(sorted_shapes) // products
    slots = [sorted_shapes[i * per_slot:(i + 1) * per_slot] for i in range(products)]
    remainder = sorted_shapes[products * per_slot:]
    if remainder:
        slots[-1].extend(remainder)
    return slots



def _center_content_a4(slide, slide_width: int) -> None:
    """Centra el contenido en plantillas de 1 producto por hoja.

    Solo se aplica cuando el slide NO tiene GroupShapes: esas plantillas tienen
    un layout intencional (precio a la izquierda, caja banco a la derecha) y
    estirar todos los shapes causaría solapamientos.
    """
    if any(hasattr(s, 'shapes') for s in slide.shapes):
        return
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        shape.left = 0
        shape.width = slide_width
        for para in shape.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER


def _align_bank_group_a4(slide) -> None:
    """Baja el GroupShape de oferta bancaria para alinearlo con el borde inferior del precio."""
    price_shape = None
    group_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'shapes'):
            group_shape = shape
        elif shape.has_text_frame:
            t = _shape_text(shape)
            if re.search(r"<<Precio\d*>>", t) or re.search(r"Precio\s+\d+", t):
                price_shape = shape
    if price_shape is not None and group_shape is not None:
        group_shape.top = price_shape.top + price_shape.height - group_shape.height


def _add_slide_from_template(prs, layout, template_shape_xmls):
    new_slide = prs.slides.add_slide(layout)
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag not in ("nvGrpSpPr", "grpSpPr"):
            sp_tree.remove(child)
    for xml_elem in template_shape_xmls:
        sp_tree.append(copy.deepcopy(xml_elem))
    return new_slide


# ---------------------------------------------------------------------------
# Main generation function
# ---------------------------------------------------------------------------

def generate_pptx_bytes(
    excel_bytes: bytes,
    template_bytes: bytes,
    vigencia: str,
    aclaracion: str,
    otra_alcohol: str,
    banco: str = "",
) -> bytes:
    products = load_products_from_bytes(excel_bytes, vigencia, aclaracion, otra_alcohol, banco)

    prs = Presentation(io.BytesIO(template_bytes))
    if not prs.slides:
        raise ValueError("La plantilla PPTX está vacía.")

    # Single-slide templates: slide 0 is the product template.
    # Multi-slide templates: slide 0 is a cover kept as-is, slide 1 is the template.
    template_slide = prs.slides[0] if len(prs.slides) == 1 else prs.slides[1]
    layout = template_slide.slide_layout

    template_shape_xmls = [copy.deepcopy(shape._element) for shape in template_slide.shapes]

    # Detect products per slide from the template's slot structure
    initial_slots = _get_slots(list(template_slide.shapes))
    products_per_slide = max(len(initial_slots), 1)

    groups = [products[i:i + products_per_slide] for i in range(0, len(products), products_per_slide)]

    for idx, group in enumerate(groups):
        slide = template_slide if idx == 0 else _add_slide_from_template(prs, layout, template_shape_xmls)
        cur_slots = _get_slots(list(slide.shapes))
        for i, product in enumerate(group):
            if i < len(cur_slots):
                _fill_slot(cur_slots[i], product, adjust_p1=(products_per_slide == 1))
        for i in range(len(group), products_per_slide):
            if i < len(cur_slots):
                _clear_slot(cur_slots[i])
        if products_per_slide == 1:
            _center_content_a4(slide, prs.slide_width)
            _align_bank_group_a4(slide)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Excel template generation
# ---------------------------------------------------------------------------

def generate_template_bytes() -> bytes:
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.utils import get_column_letter
    from openpyxl.worksheet.datavalidation import DataValidation

    HEADERS = ["Categoria", "subcategoria", "OFERTADET", "DESCRIPCION", "PRECIO", "OFERTA", "MONEDA"]
    EXAMPLES: list[tuple] = [
        ("ALIMENTOS",           "GALLETITAS",        "Precio fijo",  "Galletitas OREO 117g",          1500,  "",       "$"),
        ("BEBIDAS SIN ALCOHOL", "GASEOSAS",           "Combo",        "Coca-Cola 2.25L",               2500,  "2X4500", "$"),
        ("BEBIDAS SIN ALCOHOL", "AGUA",               "M x N",        "Agua SALUS 1.5L",               800,   "",       "$"),
        ("LIMPIEZA",            "LIMPIADORES",        "% descuento",  "Lavandina AYUDIN 2L",           850,   "",       "$"),
        ("FIAMBRES Y QUESOS",   "QUESOS",             "Precio fijo",  "Queso Barra por Kg.",           12000, "",       "$"),
        ("BEBIDAS CON ALCOHOL", "VINOS",              "Precio fijo",  "Vino NORTON Malbec 750ml",      3200,  "",       "$"),
        ("ELECTRODOMESTICOS",   "ELECTRODOMESTICOS",  "Precio fijo",  "Licuadora PHILIPS HR2100",      45,    "",       "U$S"),
    ]

    wb = openpyxl.Workbook()

    # ── Cenefas sheet ──────────────────────────────────────────────────────
    ws = wb.active
    ws.title = "Cenefas"

    header_fill = PatternFill("solid", fgColor="1E3A5F")
    header_font = Font(bold=True, color="FFFFFF", size=11)
    even_fill   = PatternFill("solid", fgColor="EEF2F7")

    for col, name in enumerate(HEADERS, 1):
        cell = ws.cell(row=1, column=col, value=name)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center", vertical="center")

    for row_idx, data in enumerate(EXAMPLES, 2):
        fill = even_fill if row_idx % 2 == 0 else None
        for col_idx, value in enumerate(data, 1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            cell.alignment = Alignment(vertical="center")
            if fill:
                cell.fill = fill

    col_widths = [22, 20, 14, 36, 10, 12, 10]
    for col, width in enumerate(col_widths, 1):
        ws.column_dimensions[get_column_letter(col)].width = width
    ws.row_dimensions[1].height = 26
    ws.freeze_panes = "A2"

    dv_tipo = DataValidation(type="list", formula1='"Precio fijo,% descuento,Combo,M x N"', allow_blank=True)
    dv_tipo.sqref = "C2:C5000"
    ws.add_data_validation(dv_tipo)

    dv_moneda = DataValidation(type="list", formula1='"$,U$S"', allow_blank=True)
    dv_moneda.sqref = "G2:G5000"
    ws.add_data_validation(dv_moneda)

    # ── Instrucciones sheet ────────────────────────────────────────────────
    ws2 = wb.create_sheet("Instrucciones")
    ws2.column_dimensions["A"].width = 20
    ws2.column_dimensions["B"].width = 45
    ws2.column_dimensions["C"].width = 32
    ws2.column_dimensions["D"].width = 35

    inst_header_font = Font(bold=True, color="FFFFFF", size=11)
    inst_header_fill = PatternFill("solid", fgColor="1E3A5F")
    inst_even_fill   = PatternFill("solid", fgColor="EEF2F7")

    inst_cols = ["Columna", "Descripción", "Valores aceptados", "Notas"]
    for col, name in enumerate(inst_cols, 1):
        cell = ws2.cell(row=1, column=col, value=name)
        cell.fill = inst_header_fill
        cell.font = inst_header_font
        cell.alignment = Alignment(horizontal="center", vertical="center")
    ws2.row_dimensions[1].height = 24

    rows = [
        ("Categoria",    "Categoría del producto",              "Texto libre",                     "Ej: ALIMENTOS, BEBIDAS CON ALCOHOL, CARNES"),
        ("subcategoria", "Subcategoría del producto",           "Texto libre",                     "QUESOS y FIAMBRES: precio por 100g si desc tiene 'kg'. CARNES/FIAMBRES/QUESOS: sin unidad."),
        ("OFERTADET",    "Tipo de oferta",                      "Precio fijo / % descuento / Combo / M x N", "Determina cómo se muestra el precio en la cenefa"),
        ("DESCRIPCION",  "Nombre del producto tal como aparece","Texto libre",                     "Las palabras en MAYÚSCULAS se muestran en negrita"),
        ("PRECIO",       "Precio unitario (número)",            "Número (ej: 1500, 45.90)",        "Para Combo: precio individual. Para dólares usa MONEDA=U$S"),
        ("OFERTA",       "Solo para Combo: cantidad y precio total", "Formato: 2X4500",            "2X4500 = 2 unidades por $4500. Vacío para otros tipos."),
        ("MONEDA",       "Moneda del precio",                   "$ o U$S",                         "$ = pesos uruguayos. U$S = dólares."),
    ]
    for row_idx, data in enumerate(rows, 2):
        fill = inst_even_fill if row_idx % 2 == 0 else None
        for col_idx, value in enumerate(data, 1):
            cell = ws2.cell(row=row_idx, column=col_idx, value=value)
            cell.alignment = Alignment(vertical="center", wrap_text=True)
            if fill:
                cell.fill = fill
        ws2.row_dimensions[row_idx].height = 36

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()
