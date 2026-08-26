"""Renderer de componentes v2 — genera PPTX desde definición JSON de componentes."""
import copy
import io
import math
import re

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt

from app.services.cenefas.data_engine import load_products_from_bytes
from app.services.cenefas.font_metrics import ancho_texto_cm
from app.services.cenefas.formatters import split_caps
from app.services.cenefas.layout_engine import compute_layout, get_format
from app.services.cenefas.rules_engine import apply_visibility, evaluate_rules
from app.services.cenefas.variables import DECIMAL_OF, PRICE_VARS

# ---------------------------------------------------------------------------
# Dimensiones de slide por formato
# ---------------------------------------------------------------------------

FORMAT_SLIDES: dict[str, tuple] = {
    "a4":      (Cm(21.0),  Cm(29.7)),
    "a3":      (Cm(29.7),  Cm(42.0)),
    "3xa4":    (Cm(21.0),  Cm(29.7)),   # A4 portrait completo, 3 franjas verticales
    "pinchos": (Cm(21.0),  Cm(29.7)),   # A4 portrait completo, grilla 3×2
    "a5":      (Cm(14.85), Cm(21.0)),
    "6xa4":    (Cm(21.0),  Cm(29.7)),   # A4 portrait completo, grilla 3×2 (arte propio)
}

ALIGN_MAP = {
    "left":   PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right":  PP_ALIGN.RIGHT,
}

# ---------------------------------------------------------------------------
# Transforms de valores
# ---------------------------------------------------------------------------

def apply_transform(value: str, transform: str | None) -> str:
    """Aplica la transformación al valor del campo."""
    if not value or not transform or transform in ("none", "smart_bold"):
        return value or ""

    if transform in ("price_full", "price_integer", "price_decimal"):
        # value ya viene formateado como "$1.250,90" desde data_engine
        num = re.sub(r"[^\d.,]", "", value)
        if transform == "price_full":
            return value
        if transform == "price_integer":
            return num.rsplit(",", 1)[0] if "," in num else num
        if transform == "price_decimal":
            return ("," + num.rsplit(",", 1)[1]) if "," in num else ""

    if transform == "combo_quantity":
        m = re.match(r"(\d+X)", value.upper())
        return m.group(1) if m else value

    if transform == "combo_price":
        return value  # ya es el precio formateado

    if transform == "uppercase":
        return value.upper()

    return value


# ---------------------------------------------------------------------------
# Achique automático de texto que no entra
# ---------------------------------------------------------------------------
#
# Única excepción a la regla de "el motor respeta el PPTX tal cual". Aplica a
# la descripción y a los precios: los nombres reales de gestión pasan seguido
# de los 60 caracteres y desbordan sobre el precio de al lado, y un precio de
# cuatro dígitos ("1.919") no entra en un cuadro calibrado para tres --
# PowerPoint lo parte al medio ("1.91" + "9", visto en un cartel real).
#
# Se sacó en 08/2026 junto con el resto de los ajustes automáticos y volvió a
# pedido explícito, primero para la descripción y después para el precio.
#
# Lo que NO volvió: el corrimiento vertical del cuadro de precio cuando la
# descripción desborda. Acá solo se achica texto; ninguna caja se mueve de
# donde la puso el diseño.


# Inset interno por defecto de PowerPoint: 0,1" a cada lado (lIns/rIns =
# 91440 EMU). Antes se descontaba 0,4 cm y el word-wrap simulado cortaba una
# palabra más tarde que el real.
_INSET_CM = 0.508

# Interlineado tipico de una caja sin espaciado explicito.
_INTERLINEADO = 1.2

# Letras que bajan de la línea de base.
_DESCENDENTES = frozenset("gjpqy")


def _alto_ultima_linea(texto: str) -> float:
    """Alto de la última línea, en múltiplos del tamaño de fuente.

    El interlineado separa una línea de la siguiente, pero debajo de la última
    no hay nada que separar: lo que puede chocar con el cuadro de abajo es la
    TINTA, y hasta dónde llega depende de qué diga el texto.

    Contando 1,2 em también para la última línea, un precio de una sola línea a
    140 pt "necesitaba" 5,93 cm donde el diseño le da 5,51 -- y el motor
    terminaba achicando los 145 precios de la A5 uno por uno sin que ninguno
    estuviera pisando nada. Un precio son dígitos: no bajan de la base, su tinta
    no pasa de la altura de mayúscula.
    """
    if _DESCENDENTES & set(texto):
        return 1.15
    if any(c.islower() for c in texto):
        return 1.05
    return 0.95


def _alto_texto_cm(lineas: int, font_size: float, texto: str) -> float:
    """Alto que ocupa la tinta de un texto de N líneas a ese tamaño."""
    factor = (lineas - 1) * _INTERLINEADO + _alto_ultima_linea(texto)
    return factor * font_size / 72 * 2.54


def _ancho_medido_cm(
    texto: str, font_size: float, font_family: str | None, bold: bool
) -> float:
    """Ancho del texto tal como se va a dibujar, en centímetros.

    Se mide parte por parte porque el renderer pone en negrita las palabras en
    mayúsculas --la marca: "SER", "LA SERENÍSIMA"-- vía split_caps. Medir todo
    el texto con un único flag de negrita sobra o falta ancho según el caso, y
    en cajas ajustadas eso alcanza para errar por una línea entera, que es
    justo la diferencia entre "entra" y "se monta sobre el precio".
    """
    return sum(
        ancho_texto_cm(parte, font_size, font_family, bold or es_mayus)
        for parte, es_mayus in split_caps(texto)
    )


def _estimate_wrapped_lines(
    text: str, box_width_cm: float | None, font_size: float | None,
    bold: bool = False, font_family: str | None = None,
) -> int:
    """A cuántas líneas se parte el texto con word-wrap en una caja de ese ancho.

    Simula el criterio "voraz" de PowerPoint (agregar palabras hasta que no
    entran más) midiendo cada palabra con las métricas REALES de la tipografía
    (ver font_metrics.py). Antes se usaba un ancho de caracter promedio
    inventado y erraba de a una línea entera, que es justo la diferencia entre
    "entra" y "se le monta encima al precio".
    """
    if not text or not box_width_cm or not font_size:
        return 1
    usable_cm = max(0.1, box_width_cm - _INSET_CM)

    lineas = 1
    actual = ""
    for palabra in text.split():
        tentativa = palabra if not actual else actual + " " + palabra
        if _ancho_medido_cm(tentativa, font_size, font_family, bold) > usable_cm and actual:
            lineas += 1
            actual = palabra
        else:
            actual = tentativa
    return lineas


def _texto_resuelto(comp: dict, product: dict) -> str:
    """El texto que ese componente va a imprimir, ya con los valores puestos."""
    segs = comp.get("segments")
    if segs:
        partes = []
        for seg in segs:
            if seg.get("type") == "variable":
                partes.append(str(product.get(seg.get("value"), "") or ""))
            else:
                partes.append(str(seg.get("value", "") or ""))
        return "".join(partes)
    variable = comp.get("variable")
    if variable:
        return str(product.get(variable, "") or "")
    return str(comp.get("static_value", "") or "")


# Piso de achique: por debajo de esto el texto deja de ser legible en un cartel
# de góndola, y es preferible que se note el desborde a imprimir algo que nadie
# puede leer de lejos.
_FIT_MIN_SCALE = 0.55

# Se achica cualquier cuadro que traiga DATO del Excel y no entre. Un cuadro de
# texto fijo del diseño ("OFERTA", "PRECIO REGULAR") no se toca nunca: su
# contenido no cambia entre productos, así que si el diseñador lo dejó justo,
# está justo a propósito.
#
# Los decimales quedan afuera: son siempre dos dígitos, nunca desbordan, y
# achicarlos desalinearía la coma respecto del entero de al lado.
#
# Antes la lista era sólo {descripción + precios} y por eso un código de varios
# SKU ("594879/80/81/82/83 -593838/39/40 - 621032 - ...", 86 caracteres) se
# partía en tres líneas y se montaba sobre la descripción.
_FIT_EXCLUIDAS: frozenset = frozenset(DECIMAL_OF.values())


def _variables_del_componente(c: dict) -> set[str]:
    """Las variables que ese componente imprime, sea directo o por segmentos."""
    if c.get("variable"):
        return {c["variable"]}
    return {
        seg["value"] for seg in (c.get("segments") or [])
        if seg.get("type") == "variable" and seg.get("value")
    }


def _segmentos_medibles(comp: dict, product: dict, escala: float = 1.0) -> list[tuple[str, float]]:
    """(texto, tamaño) de cada pedazo del cuadro, ya resuelto contra el producto.

    Un cuadro de precio tiene tres pedazos con tamaños MUY distintos: el "$" a
    100 pt, el número a 180 y los centavos a 90. Medirlo todo con un solo
    tamaño --el del primer run, que es el "$"-- daba un ancho como la mitad del
    real: el motor creía que "$147,20" entraba en 12,35 cm, no achicaba, y
    PowerPoint terminaba partiendo los centavos en dos líneas que se caían
    sobre la cenefa de abajo (visto en la 3xA4 del 27/08).

    Devuelve [] cuando el cuadro no tiene tamaños por segmento; ahí la medición
    de siempre, con un único tamaño, es correcta.
    """
    segs = comp.get("segments") or []
    if not any((seg.get("style") or {}).get("font_size") for seg in segs):
        return []
    # 18 pt es el default real de PowerPoint cuando el cuadro no declara tamano.
    base = (comp.get("style") or {}).get("font_size") or 18.0
    salida: list[tuple[str, float]] = []
    for seg in segs:
        if seg.get("type") == "variable":
            texto = str(product.get(seg.get("value"), "") or "")
        else:
            texto = str(seg.get("value", "") or "")
        if not texto:
            continue
        salida.append((texto, ((seg.get("style") or {}).get("font_size") or base) * escala))
    return salida


def _entra_por_segmentos(
    piezas: list[tuple[str, float]], box_width_cm: float, box_height_cm: float | None,
    bold: bool, font_family: str | None,
) -> bool:
    """Igual que _entra_en_caja pero sumando el ancho pedazo por pedazo.

    Un precio no tiene espacios, así que no hay dónde cortarlo por palabra: o
    entra entero en una línea o PowerPoint lo parte al medio. Por eso se exige
    que la suma de los anchos entre en el ancho útil, sin word-wrap posible.
    """
    usable_cm = max(0.1, box_width_cm - _INSET_CM)
    ancho = sum(_ancho_medido_cm(t, sz, font_family, bold) for t, sz in piezas)
    if ancho > usable_cm:
        return False
    if box_height_cm:
        mayor = max(sz for _, sz in piezas)
        if _alto_texto_cm(1, mayor, "".join(t for t, _ in piezas)) > box_height_cm:
            return False
    return True


def _entra_en_caja(
    texto: str, box_width_cm: float, box_height_cm: float | None,
    font_size: float, bold: bool, font_family: str | None,
) -> bool:
    """True si el texto entra en el cuadro a ese tamaño de fuente.

    Dos condiciones, las dos obligatorias:

    1. **Ancho**: la palabra más larga tiene que entrar en una línea. Es el
       chequeo que salva a los precios: "1.919" no tiene espacios, así que no
       hay dónde cortarlo por palabra -- PowerPoint lo parte al medio
       ("1.91" + "9", visto en un cartel real) y el chequeo de alto ni se
       entera, porque para él sigue siendo una sola línea.
    2. **Alto**: las líneas que resultan del word-wrap tienen que entrar en el
       alto disponible. Es el chequeo que salva a las descripciones largas.
    """
    usable_cm = max(0.1, box_width_cm - _INSET_CM)

    palabra_larga = max(texto.split() or [texto], key=len)
    if _ancho_medido_cm(palabra_larga, font_size, font_family, bold) > usable_cm:
        return False

    if box_height_cm:
        lineas = _estimate_wrapped_lines(texto, box_width_cm, font_size, bold, font_family)
        if _alto_texto_cm(lineas, font_size, texto) > box_height_cm:
            return False

    return True


def _fit_font_size(
    texto: str, box_width_cm: float | None, box_height_cm: float | None,
    base_font_size: float | None, bold: bool = False, font_family: str | None = None,
) -> float | None:
    """El tamaño de fuente MÁS GRANDE al que el texto entra en su cuadro.

    Se busca por bisección en vez de calcular una escala de una sola pasada.
    La escala directa (alto_disponible / alto_necesario) achica de más: al
    bajar el tamaño el texto pasa a ocupar menos líneas, así que el alto que
    hacía falta era mucho menor que el que se usó para calcularla. En un caso
    real la descripción del 3xA4 terminaba en 18,8 pt cuando en 22,5 pt ya
    entraba. Achicar lo menos posible es respetar el diseño.

    "Entra" es monótono respecto del tamaño --si entra a N pt, entra a
    cualquier tamaño menor-- que es lo que hace válida la bisección.
    """
    if not texto or not box_width_cm or not base_font_size:
        return base_font_size

    if _entra_en_caja(texto, box_width_cm, box_height_cm, base_font_size, bold, font_family):
        return base_font_size

    # Piso de legibilidad: por debajo de esto es preferible que se note el
    # desborde antes que imprimir algo ilegible a dos metros de distancia.
    minimo = base_font_size * _FIT_MIN_SCALE
    mejor = minimo
    lo, hi = minimo, base_font_size
    for _ in range(12):
        medio = (lo + hi) / 2.0
        if _entra_en_caja(texto, box_width_cm, box_height_cm, medio, bold, font_family):
            mejor = medio
            lo = medio
        else:
            hi = medio

    # Se redondea hacia ABAJO a medio punto: hacia arriba podría dejar de entrar.
    return max(minimo, math.floor(mejor * 2.0) / 2.0)


# Pares de variables que NUNCA se dibujan juntas: si la primera trae valor, la
# segunda no se dibuja aunque el diseno tenga su cuadro.
#
# `promoOferta` (el literal de la mecanica, "6x4") y `precioOferta` (el precio)
# ocupan EL MISMO lugar del cartel en las plantillas de Redexpres: el diseno
# pone el cuadro de promoOferta encima del del precio. Superponer dos cuadros de
# texto no tapa nada --se leen los dos encimados-- asi que la exclusion se
# resuelve acá: cuando hay mecanica se dibuja el literal, cuando no, el precio.
#
# Los decimales siguen a su entero: sin esto, un M x N de "6x4" imprimia el
# literal y al lado le quedaba colgado el ",33" del precio que ya no se ve.
_EXCLUYENTES: dict[str, tuple[str, ...]] = {
    "promoOferta": ("precioOferta", "decimalPrecioOferta"),
}


# Cuanto tiene que bajar otro cuadro para contar como "el de abajo" y no como
# un vecino puesto a la misma altura. Medio centimetro: menos que eso, en un
# diseno hecho a mano, es desprolijidad de posicionamiento, no una fila nueva.
_SEPARACION_MIN_CM = 0.5


def _alto_disponible_cm(comp: dict, comps: list[dict]) -> float | None:
    """Cuánto puede crecer hacia abajo un cuadro antes de pisar a otro.

    El alto declarado de la caja NO sirve como límite: los diseñadores la
    dibujan del alto de UNA línea aunque abajo haya lugar de sobra, y el texto
    simplemente se desborda por fuera. Usarlo hacía que cualquier descripción
    de dos líneas se achicara de golpe cuando en el cartel entraba perfecta.

    El límite real es el cuadro de abajo. Solo cuentan los que están realmente
    debajo y no al costado: se exige que se solapen horizontalmente en más de
    la mitad del ancho. Sin ese filtro, el cuadro del decimal --que va pegado
    al precio, apenas más abajo-- se tomaría como techo del precio y lo
    achicaría a la nada.
    """
    b = comp.get("computed_bounds") or comp.get("base_bounds") or {}
    y, h, x, w = b.get("y"), b.get("height"), b.get("x"), b.get("width")
    if y is None or h is None or x is None or w is None:
        return h

    techo = None
    for otro in comps:
        if otro is comp:
            continue
        ob = otro.get("computed_bounds") or otro.get("base_bounds") or {}
        oy, ox, ow = ob.get("y"), ob.get("x"), ob.get("width")
        if oy is None or ox is None or ow is None:
            continue
        # Un cuadro que arranca a la MISMA altura que este no esta abajo: esta
        # al lado. Sin este margen, el cuadro de "Comprando 6" que en una franja
        # quedo 0,9 mm por debajo del techo del precio se tomaba como el piso
        # del precio: alto disponible 0,09 cm, y el precio se achicaba al minimo
        # -- visible como una franja con el precio la mitad de grande que las
        # otras dos, con el mismo diseno.
        if oy <= y + _SEPARACION_MIN_CM:
            continue
        # Hay que distinguir "abajo" de "al costado", y el ancho del solape solo
        # no alcanza. Dos casos reales que se parecen y necesitan lo contrario:
        #
        #   descripcion 4,16 ancho 15,46  /  precio tachado 4,16 ancho 7,00
        #       -> solape 7,00. Es el renglon de ABAJO: la descripcion no puede
        #          crecer sobre el. Con el filtro de "la mitad de mi ancho"
        #          (7,73) no contaba, y una descripcion de cuatro renglones se
        #          imprimia encima del tachado.
        #
        #   precio 11,29 ancho 12,35  /  cuadro del decimal 18,14 ancho 2,54
        #       -> solape 2,54. NO es el renglon de abajo: son los centavos, AL
        #          COSTADO del numero. Tomarlo como techo achica el precio a la
        #          nada.
        #
        # Lo que los separa no es cuanto se solapan sino DONDE arranca el otro:
        # un cuadro alineado con nuestro borde izquierdo es la linea siguiente;
        # uno que arranca pasada la mitad de nuestro ancho es una anotacion al
        # costado.
        solape = min(x + w, ox + ow) - max(x, ox)
        if solape <= 0:
            continue
        tapa_lo_ancho = solape >= w * 0.5
        arranca_donde_yo = abs(ox - x) <= w * 0.15
        if not (tapa_lo_ancho or arranca_donde_yo):
            continue
        techo = oy if techo is None else min(techo, oy)

    if techo is None:
        return h
    # El hueco manda, aunque sea MENOR que el alto dibujado de la caja. Los
    # diseños reales tienen la caja de la descripción solapada con la de
    # abajo (en la A5: caja de 3,68 cm pero solo 2,80 cm hasta "PRECIO
    # REGULAR"), así que quedarse con el mayor de los dos --como se hacía
    # antes-- era justamente ignorar la colisión que hay que evitar.
    return techo - y


def _ancho_util_cm(bounds: dict, ancho_pagina_cm: float | None) -> float | None:
    """Cuanto ancho tiene REALMENTE el cuadro, sin pasarse del borde del papel.

    Un cuadro puede estar declarado mas ancho que lo que queda de hoja: el
    autoajuste de PowerPoint estira la caja del precio hasta 12,36 cm arrancando
    en x=11,28, o sea 2,64 cm afuera de una A4. Midiendo contra la caja el motor
    cree que "$1.100" entra, no achica nada, y el precio se imprime cortado por
    el borde. El limite verdadero es el que llegue primero: el ancho de la caja
    o lo que queda de papel.
    """
    ancho = bounds.get("width")
    if not ancho or not ancho_pagina_cm:
        return ancho
    x = bounds.get("x")
    if x is None:
        return ancho
    return max(0.5, min(ancho, ancho_pagina_cm - x))


def _fit_text_to_box(
    comps: list[dict], product: dict, ancho_pagina_cm: float | None = None,
) -> list[dict]:
    """Achica la fuente de los cuadros cuyo valor no entra.

    Aplica a la descripción y a los precios: son cuadros de texto
    independientes apilados, sin auto-layout compartido, así que el texto de
    uno se desborda sobre otro (o se parte al medio) si no se achica.

    La idea es RESPETAR el tamaño del diseño y tocarlo solo cuando de verdad
    no entra. Por eso el límite vertical no es el alto dibujado de la caja
    sino el espacio hasta el cuadro de abajo, y el ancho del texto se mide con
    la tipografía real y no con un promedio.

    Solo cambia tamaños de fuente. Ninguna caja se mueve de donde la puso el
    diseño.
    """
    result = []
    for c in comps:
        usadas = _variables_del_componente(c)
        if c.get("type") != "text" or not usadas or usadas <= _FIT_EXCLUIDAS:
            result.append(c)
            continue

        style = c.get("style", {})
        base_font_size = style.get("font_size")
        bold = bool(style.get("font_bold"))
        familia = style.get("font_family")
        bounds = c.get("computed_bounds") or c.get("base_bounds") or {}
        texto = _texto_resuelto(c, product)

        piezas = _segmentos_medibles(c, product)
        if piezas:
            # Cuadro con tamaños mezclados (el precio): se busca la escala más
            # grande a la que la suma de los pedazos entra, y se aplica a todos
            # por igual para no desalinear la coma con el entero.
            ancho_caja = _ancho_util_cm(bounds, ancho_pagina_cm) or 0
            alto_caja = _alto_disponible_cm(c, comps)
            escala = 1.0
            if ancho_caja and not _entra_por_segmentos(piezas, ancho_caja, alto_caja, bold, familia):
                lo, hi = _FIT_MIN_SCALE, 1.0
                escala = _FIT_MIN_SCALE
                for _ in range(12):
                    medio = (lo + hi) / 2.0
                    if _entra_por_segmentos(_segmentos_medibles(c, product, medio),
                                            ancho_caja, alto_caja, bold, familia):
                        escala, lo = medio, medio
                    else:
                        hi = medio
            fitted = round(base_font_size * escala, 1) if base_font_size else base_font_size
        else:
            fitted = _fit_font_size(
                texto, _ancho_util_cm(bounds, ancho_pagina_cm), _alto_disponible_cm(c, comps),
                base_font_size, bold, familia,
            )
        if fitted == base_font_size:
            result.append(c)
            continue

        nuevo_style = {**style, "font_size": fitted}
        # El alto de línea acompaña al achique. Es el run vacío que
        # _populate_text_frame agrega para fijar el interlineado: si queda en el
        # tamaño original mientras el número baja, el renglón conserva el alto
        # de antes y el texto se planta más abajo de donde lo puso el diseño.
        if style.get("line_height_pt") and base_font_size:
            nuevo_style["line_height_pt"] = round(
                style["line_height_pt"] * (fitted / base_font_size), 1)
        c = {**c, "style": nuevo_style}
        # En un componente multi-segmento cada segmento lleva su propio
        # font_size, que pisa al del componente en _populate_text_frame. Sin
        # esto el achique se descartaba en silencio para cualquier cuadro
        # importado como multi-segmento (que ahora son casi todos).
        segs = c.get("segments")
        if segs:
            escala = fitted / base_font_size
            c["segments"] = [
                {**seg, "style": {**seg["style"], "font_size": seg["style"]["font_size"] * escala}}
                if seg.get("style", {}).get("font_size") else seg
                for seg in segs
            ]
        result.append(c)
    return result


def hex_to_rgb(hex_color: str | None) -> RGBColor:
    if not hex_color:
        return RGBColor(0x1E, 0x29, 0x3B)
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    try:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except ValueError:
        return RGBColor(0x1E, 0x29, 0x3B)


# ---------------------------------------------------------------------------
# Renderizado de componentes individuales
# ---------------------------------------------------------------------------

def _enable_normAutofit(tf) -> None:
    body_pr = tf._txBody.find(qn("a:bodyPr"))
    if body_pr is None:
        return
    for tag in (qn("a:spAutoFit"), qn("a:noAutofit")):
        el = body_pr.find(tag)
        if el is not None:
            body_pr.remove(el)
    if body_pr.find(qn("a:normAutofit")) is None:
        body_pr.append(etree.Element(qn("a:normAutofit")))


def _populate_text_frame(tf, comp: dict, value: str) -> None:
    """Arma el contenido de un text_frame ya existente — separado de
    add_text_component para poder reusarlo tanto al crear una caja de texto
    nueva como al reescribir el texto de un shape original preservado
    (ver _place_component)."""
    style    = comp.get("style", {})
    segments = comp.get("segments")

    tf.word_wrap = True

    # Vertical anchor (preserved from original PPTX)
    vertical_align = style.get("vertical_align")
    if vertical_align:
        try:
            body_pr = tf._txBody.find(qn("a:bodyPr"))
            if body_pr is not None:
                body_pr.set("anchor", vertical_align)
        except Exception:
            pass

    if style.get("auto_fit", False):
        _enable_normAutofit(tf)

    tf.clear()  # saca cualquier texto previo (ej. "<<Descripción>>" del original)
    transform = comp.get("transform", "none")
    p = tf.paragraphs[0]
    p.alignment = ALIGN_MAP.get(style.get("align", "center"), PP_ALIGN.CENTER)

    if segments:
        # Multi-segment: each segment gets its own run with per-segment style overrides.
        # Variable segments have their value pre-resolved as "_resolved" by _render_slide.
        for seg in segments:
            seg_val = seg.get("_resolved", seg.get("value", ""))
            if not seg_val:
                continue
            seg_style = {**style}
            if seg.get("style"):
                seg_style.update(seg["style"])
            seg_transform = seg.get("transform") or "none"
            if seg_transform == "smart_bold":
                for part, is_bold in split_caps(seg_val):
                    if part:
                        run = p.add_run()
                        run.text = part
                        _apply_run_style(run, seg_style, bold_override=is_bold)
            else:
                if seg_transform not in (None, "none"):
                    seg_val = apply_transform(seg_val, seg_transform)
                run = p.add_run()
                run.text = seg_val
                _apply_run_style(run, seg_style)
    elif transform == "smart_bold":
        for segment, is_bold in split_caps(value):
            if not segment:
                continue
            run = p.add_run()
            run.text = segment
            _apply_run_style(run, style, bold_override=is_bold)
    else:
        run_style = style
        run = p.add_run()
        run.text = value
        _apply_run_style(run, run_style)

    # Replicate empty spacer run used in original PPTX to set a larger line height.
    # Without this, anchor=b positions text much lower than the original.
    line_height_pt = style.get("line_height_pt")
    if line_height_pt and line_height_pt != style.get("font_size"):
        spacer = p.add_run()
        spacer.text = ""
        spacer.font.size = Pt(line_height_pt)


def add_text_component(slide, comp: dict, value: str) -> None:
    bounds = comp["computed_bounds"]
    txBox = slide.shapes.add_textbox(
        Cm(bounds["x"]), Cm(bounds["y"]),
        Cm(max(bounds["width"],  0.5)),
        Cm(max(bounds["height"], 0.3)),
    )
    _populate_text_frame(txBox.text_frame, comp, value)


def _apply_run_style(run, style: dict, bold_override: bool | None = None) -> None:
    font = run.font
    if style.get("font_size"):
        font.size = Pt(style["font_size"])
    font.bold = bold_override if bold_override is not None else style.get("font_bold", False)
    if style.get("color"):
        font.color.rgb = hex_to_rgb(style["color"])
    if style.get("font_family"):
        font.name = style["font_family"]
    # python-pptx no expone tachado en Font — hay que bajar al XML crudo.
    if style.get("strikethrough"):
        run._r.get_or_add_rPr().set("strike", "sngStrike")
    # Tampoco expone la voladita. Es lo que mantiene el "$" y los centavos
    # arriba de la línea de base del número grande.
    if style.get("baseline"):
        run._r.get_or_add_rPr().set("baseline", str(style["baseline"]))


def add_shape_component(slide, comp: dict) -> None:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    bounds = comp["computed_bounds"]
    style  = comp.get("style", {})

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Cm(bounds["x"]), Cm(bounds["y"]),
        Cm(max(bounds["width"],  0.5)),
        Cm(max(bounds["height"], 0.3)),
    )
    # Sin borde visible
    shape.line.fill.background()

    if style.get("background_color"):
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(style["background_color"])
    else:
        shape.fill.background()


_WMF_EXTS = {"wmf", "emf"}

def add_image_from_data(slide, comp: dict) -> None:
    """Embebe una imagen base64 en el slide.
    Formatos web (JPEG/PNG): usa add_picture normal.
    Formatos vectoriales (WMF/EMF): embebe via XML directo, sin PIL."""
    import base64 as _b64
    bounds   = comp["computed_bounds"]
    img_ext  = (comp.get("image_ext") or "").lower()

    try:
        img_bytes = _b64.b64decode(comp["image_data"])
    except Exception:
        add_image_placeholder(slide, comp, comp.get("name", "imagen"))
        return

    if img_ext in _WMF_EXTS:
        try:
            _embed_vector_image(slide, img_bytes, img_ext, bounds)
        except Exception:
            # Only show placeholder for variable images (product images).
            # Decorative static images (no variable) are silently skipped.
            if comp.get("variable"):
                add_image_placeholder(slide, comp, comp.get("name", "imagen"))
    else:
        try:
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Cm(bounds["x"]),
                Cm(bounds["y"]),
                Cm(max(bounds["width"],  0.1)),
                Cm(max(bounds["height"], 0.1)),
            )
        except Exception:
            add_image_placeholder(slide, comp, comp.get("name", "imagen"))


def _embed_vector_image(slide, img_bytes: bytes, ext: str, bounds: dict) -> None:
    """Embebe WMF/EMF directamente en el XML del slide sin pasar por PIL."""
    import hashlib
    from pptx.opc.part import Part
    from pptx.opc.packuri import PackURI

    content_types = {"wmf": "image/x-wmf", "emf": "image/x-emf"}
    ct  = content_types.get(ext, "image/x-wmf")
    h   = hashlib.md5(img_bytes).hexdigest()[:12]
    uri = PackURI(f"/ppt/media/img_{h}.{ext}")

    img_part = Part(uri, ct, img_bytes)
    rId = slide.part.relate_to(
        img_part,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
    )

    x  = int(Cm(bounds["x"]))
    y  = int(Cm(bounds["y"]))
    cx = int(Cm(max(bounds["width"],  0.1)))
    cy = int(Cm(max(bounds["height"], 0.1)))
    pid = abs(hash(h)) % 8000 + 1000

    pic_xml = (
        f'<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        f' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        f'<p:nvPicPr>'
        f'<p:cNvPr id="{pid}" name="img_{h[:8]}"/>'
        f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
        f'<p:nvPr/></p:nvPicPr>'
        f'<p:blipFill><a:blip r:embed="{rId}"/>'
        f'<a:stretch><a:fillRect/></a:stretch></p:blipFill>'
        f'<p:spPr><a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
        f'</p:pic>'
    )
    slide.shapes._spTree.append(parse_xml(pic_xml.encode()))


def add_image_placeholder(slide, comp: dict, label: str) -> None:
    """Placeholder para imágenes — rectángulo gris con el nombre de la variable.

    Reemplazar por descarga real de URL cuando se soporte HTTP en el renderer.
    """
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    bounds = comp["computed_bounds"]
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Cm(bounds["x"]), Cm(bounds["y"]),
        Cm(max(bounds["width"],  0.5)),
        Cm(max(bounds["height"], 0.5)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)  # slate-200
    shape.line.fill.background()

    # Label indicativo centrado
    tf = shape.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[{label}]"
    run.font.size  = Pt(9)
    run.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)  # slate-400


# ---------------------------------------------------------------------------
# Preservación del diseño original — reusar/mutar shapes reales del PPTX
# fuente en vez de reconstruir todo sobre una presentación en blanco.
# ---------------------------------------------------------------------------

def _shape_id_map(shapes) -> dict[int, object]:
    """Mapa shape_id -> shape, recorriendo grupos recursivamente (mismo
    criterio que _flatten_shapes en pptx_importer, para que el id capturado
    al importar siga siendo encontrable acá)."""
    result: dict[int, object] = {}
    for shape in shapes:
        if hasattr(shape, "shapes"):
            result.update(_shape_id_map(shape.shapes))
        else:
            try:
                result[shape.shape_id] = shape
            except Exception:
                pass
    return result


def _duplicate_slide(prs, source_slide):
    """Clona un slide completo (layout, shapes, relaciones de imagen y fondo
    propio si tiene) — python-pptx no trae esto de fábrica. Hace falta para
    generar N páginas iguales al diseño original cuando hay más de un grupo
    de productos por generar (ver render_template_to_pptx)."""
    new_slide = prs.slides.add_slide(source_slide.slide_layout)

    # add_slide ya pudo haber agregado placeholders heredados del layout —
    # los sacamos, vamos a clonar los shapes reales del slide fuente.
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    # Fondo propio del slide (si el original no lo hereda del layout/master)
    src_bg = source_slide._element.find(qn("p:cSld")).find(qn("p:bg"))
    if src_bg is not None:
        new_cSld = new_slide._element.find(qn("p:cSld"))
        new_cSld.insert(0, copy.deepcopy(src_bg))

    # Relaciones de imagen — para que los r:embed de los shapes clonados
    # sigan resolviendo a la parte correcta en el slide nuevo.
    id_map: dict[str, str] = {}
    for rel_id, rel in source_slide.part.rels.items():
        if rel.is_external or "image" not in rel.reltype:
            continue
        id_map[rel_id] = new_slide.part.relate_to(rel.target_part, rel.reltype)

    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape._element)
        if id_map:
            for blip in new_el.iter(qn("a:blip")):
                old_rid = blip.get(qn("r:embed"))
                if old_rid and old_rid in id_map:
                    blip.set(qn("r:embed"), id_map[old_rid])
        new_slide.shapes._spTree.append(new_el)

    return new_slide


def _place_component(slide, comp: dict, value: str, shape_map: dict[int, object]) -> None:
    """Coloca un componente en el slide. Si viene de una plantilla con
    diseño original preservado y matchea un shape real del archivo fuente
    (_source_shape_id), lo muta en el lugar — reposiciona y reescribe su
    contenido — en vez de crear uno nuevo, así el resto del diseño del
    archivo (fondos, logos, bordes que el importer no haya capturado como
    componente) queda intacto. Sin match, cae al comportamiento de siempre."""
    comp_type = comp.get("type", "text")
    source_id = comp.get("_source_shape_id")
    shape = shape_map.get(source_id) if source_id is not None else None

    if shape is not None:
        bounds = comp["computed_bounds"]
        try:
            shape.left   = Cm(bounds["x"])
            shape.top    = Cm(bounds["y"])
            shape.width  = Cm(max(bounds["width"],  0.1))
            shape.height = Cm(max(bounds["height"], 0.1))
        except Exception:
            pass

        if comp_type == "text" and shape.has_text_frame:
            _populate_text_frame(shape.text_frame, comp, value)
            return
        if comp_type == "shape":
            style = comp.get("style", {})
            if style.get("background_color"):
                try:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = hex_to_rgb(style["background_color"])
                except Exception:
                    pass
            return
        if comp_type == "image":
            if comp.get("image_data"):
                # Más simple/confiable reemplazar el shape entero que mutar
                # el blob de una <p:pic> ya existente.
                shape._element.getparent().remove(shape._element)
                add_image_from_data(slide, comp)
            return
        return

    # Sin shape original que matchee (template armado en el editor, o
    # componente agregado a mano después de importar) -> comportamiento
    # de siempre: crear el shape desde cero.
    if comp_type == "text":
        add_text_component(slide, comp, value)
    elif comp_type == "shape":
        add_shape_component(slide, comp)
    elif comp_type == "image":
        if comp.get("image_data"):
            add_image_from_data(slide, comp)
        else:
            add_image_placeholder(slide, comp, comp.get("variable") or "imagen")


# ---------------------------------------------------------------------------
# Render de un slide completo
# ---------------------------------------------------------------------------

def _render_slide(
    slide,
    comp_layout: list[dict],
    product: dict,
    slot_offset_x: float = 0.0,
    slot_offset_y: float = 0.0,
    missing_vars: set | None = None,
    shape_map: dict[int, object] | None = None,
) -> None:
    shape_map = shape_map or {}
    for comp in comp_layout:
        comp_type = comp.get("type", "text")
        source_id = comp.get("_source_shape_id")

        oculto = not comp.get("visible", True)
        # Excluyentes: si la que manda del par trae valor, esta no se dibuja.
        if not oculto:
            usadas = _variables_del_componente(comp)
            for manda, tapadas in _EXCLUYENTES.items():
                if usadas & set(tapadas) and str(product.get(manda, "") or "").strip():
                    oculto = True
                    break
        # Un cuadro cuyo contenido sale SOLO de variables y todas quedaron
        # vacías no tiene nada que imprimir. Borrarle el texto no alcanza: si
        # el shape tiene relleno propio --la cocarda roja de tipoOferta-- queda
        # un rectángulo de color impreso en el cartel de un producto que no
        # tiene mecánica. Se saca el shape entero.
        if not oculto and comp_type == "text" and _variables_del_componente(comp):
            partes_fijas = any(
                str(seg.get("value", "")).strip()
                for seg in (comp.get("segments") or []) if seg.get("type") == "static"
            )
            if not partes_fijas and not _texto_resuelto(comp, product).strip():
                oculto = True

        if oculto:
            shape = shape_map.get(source_id) if source_id is not None else None
            if shape is not None:
                # El padre puede ser None si este shape ya se saco: pasa cuando
                # el PPTX trae dos shapes con el MISMO id (PowerPoint lo evita,
                # pero copiar shapes entre archivos con python-pptx no) y dos
                # componentes distintos apuntan al mismo. Antes reventaba con un
                # AttributeError a mitad del render y no se generaba nada.
                padre = shape._element.getparent()
                if padre is not None:
                    padre.remove(shape._element)
            continue

        segments = comp.get("segments") if comp_type == "text" else None

        if segments:
            # Resolve variable segments from product data, store as "_resolved"
            resolved = []
            for i, seg in enumerate(segments):
                if seg.get("type") == "variable":
                    seg_var = seg.get("value", "")
                    seg_val = str(product.get(seg_var, "") or "") if seg_var else ""
                    if seg_var and seg_var not in product and missing_vars is not None:
                        missing_vars.add(seg_var)
                    # Guarda contra el símbolo de moneda duplicado. Desde
                    # 08/2026 los precios viajan SIN "$" (el símbolo es texto
                    # fijo del diseño), así que esto normalmente no se
                    # dispara; sigue acá para el caso de un Excel cargado a
                    # mano con "$899" en una plantilla que además tiene el
                    # "$" como run propio -- sin el chequeo queda "$$899".
                    # Solo saca un símbolo repetido, nunca reformatea nada.
                    if i > 0 and segments[i - 1].get("type") == "static":
                        prev = segments[i - 1].get("value", "").strip()
                        if prev in ("U$S", "$") and seg_val.strip().startswith(prev):
                            seg_val = seg_val.strip()[len(prev):].lstrip()
                else:
                    seg_val = seg.get("value", "")
                resolved.append({**seg, "_resolved": seg_val})
            comp  = {**comp, "segments": resolved}
            value = ""  # unused when segments present
        else:
            variable     = comp.get("variable")
            static_value = comp.get("static_value", "")
            raw_value    = str(product.get(variable, "") or "") if variable else static_value
            transform    = comp.get("transform", "none")
            value        = apply_transform(raw_value, transform)

            # Collect variables that are used in the template but whose column is
            # entirely absent from the Excel (key not in product at all).
            # Empty cells produce key="" — that's valid data, not a missing column.
            if variable and variable not in product and missing_vars is not None:
                missing_vars.add(variable)

        # Offset 2D para layouts multi-slot (grilla horizontal × vertical)
        if slot_offset_x > 0 or slot_offset_y > 0:
            cb = comp["computed_bounds"].copy()
            cb["x"] = cb["x"] + slot_offset_x
            cb["y"] = cb["y"] + slot_offset_y
            comp = {**comp, "computed_bounds": cb}

        _place_component(slide, comp, value, shape_map)


# ---------------------------------------------------------------------------
# Multi-slot A4 detection
# ---------------------------------------------------------------------------

def _esquina(c: dict) -> tuple[float, float]:
    """Punto de referencia de un cuadro: su esquina superior izquierda.

    NO el centro. Los diseños reales traen cajas absurdamente altas --el
    cuadro de <<precioOferta>> de la 6xA4 mide 10,5 cm de alto para un texto
    de una línea-- y su centro cae en la fila de ABAJO. Con el centro, la
    cenefa de arriba perdía el precio y la de abajo terminaba con dos.
    El borde superior izquierdo es donde el texto realmente empieza.
    """
    b = c.get("base_bounds", {}) or {}
    return (b.get("x", 0.0), b.get("y", 0.0))


def _cortar_por_huecos(valores: list[float], n_grupos: int) -> list[float]:
    """Límites que parten `valores` en n_grupos, cortando por los huecos mayores.

    Entre dos filas de cenefas hay aire; dentro de una fila los cuadros están
    pegados. Buscar los n-1 huecos más grandes encuentra esas separaciones sin
    depender de dónde esté el ancla ni de cuán alta sea cada caja.
    """
    if n_grupos <= 1 or len(valores) < n_grupos:
        return []
    ordenados = sorted(valores)
    huecos = sorted(
        ((ordenados[i + 1] - ordenados[i], i) for i in range(len(ordenados) - 1)),
        reverse=True,
    )[: n_grupos - 1]
    return sorted((ordenados[i] + ordenados[i + 1]) / 2.0 for _, i in huecos)


def _indice_por_limites(valor: float, limites: list[float]) -> int:
    for i, lim in enumerate(limites):
        if valor < lim:
            return i
    return len(limites)


# Tolerancia al ubicar un cuadro en su celda, como fracción del paso. Los
# diseños están hechos a mano: la fila 2 de la 6xA4 arranca en 7,01 cm cuando
# el paso exacto da 6,985, y sin margen ese cuadro cae en la fila de arriba.
_MARGEN_CELDA = 0.02


def _indice_por_paso(valor: float, origen: float, paso: float, n: int) -> int:
    """En qué celda de una grilla de paso fijo cae `valor`."""
    if n <= 1 or paso <= 0:
        return 0
    return max(0, min(n - 1, int((valor - origen) / paso + _MARGEN_CELDA)))


def _paso_de_grilla(valores: list[float], n_grupos: int) -> float | None:
    """Distancia entre celdas, deducida de dónde están las anclas.

    Las anclas se agrupan por los huecos grandes (una posición por celda) y el
    paso sale de la distancia entre la primera y la última. Devuelve None si
    las anclas no se separan en exactamente n_grupos posiciones.
    """
    if n_grupos <= 1:
        return 0.0
    limites = _cortar_por_huecos(valores, n_grupos)
    if len(limites) != n_grupos - 1:
        return None
    grupos: dict[int, list[float]] = {}
    for v in valores:
        grupos.setdefault(_indice_por_limites(v, limites), []).append(v)
    if len(grupos) != n_grupos:
        return None
    centros = [sum(grupos[i]) / len(grupos[i]) for i in range(n_grupos)]
    paso = (centros[-1] - centros[0]) / (n_grupos - 1)
    return paso if paso > 0 else None


def _asignar_grilla(
    non_bg: list[dict], anclas: list[dict], n_filas: int, n_cols: int
) -> list[list[dict]] | None:
    """Reparte los componentes en una grilla de n_filas x n_cols, o None si no cierra.

    El reparto es por PASO FIJO, no por el punto medio entre anclas. Las celdas
    de una cenefa son rectángulos iguales y repetidos: el ancla marca dónde
    ARRANCA cada celda, y el contenido de esa celda se extiende hacia la
    derecha y hacia abajo hasta donde arranca la siguiente.

    Cortar por el punto medio entre anclas asumía que el contenido está
    centrado en su ancla, y no lo está. En la A5 las anclas caen en x=0,00 y
    x=15,07, así que el corte quedaba en 7,54 -- y el cuadro del decimal del
    precio de la cenefa IZQUIERDA, que vive en x=10,92, se iba a la celda de la
    derecha. Resultado visible: la cenefa de la izquierda imprimía el decimal
    del producto de la derecha ("175" del producto 1 con el ",80" del producto
    2) y encima quedaba mal plantado.
    """
    paso_x = _paso_de_grilla([_esquina(c)[0] for c in anclas], n_cols)
    if paso_x is None:
        return None
    origen_x = min(_esquina(c)[0] for c in non_bg)

    columnas: dict[int, list[dict]] = {}
    for c in non_bg:
        columnas.setdefault(_indice_por_paso(_esquina(c)[0], origen_x, paso_x, n_cols), []).append(c)
    if len(columnas) != n_cols:
        return None

    ids_ancla = {id(c) for c in anclas}
    celdas: dict[tuple[int, int], list[dict]] = {}
    for col, comps_col in columnas.items():
        ys_ancla = [_esquina(c)[1] for c in comps_col if id(c) in ids_ancla]
        if len(ys_ancla) != n_filas:
            return None
        # Las filas se resuelven DENTRO de cada columna: cada columna puede
        # tener su propio corrimiento vertical (en la 6xA4 la columna derecha
        # arranca 2 mm más abajo que la izquierda).
        paso_y = _paso_de_grilla(ys_ancla, n_filas)
        if paso_y is None:
            return None
        origen_y = min(_esquina(c)[1] for c in comps_col)
        for c in comps_col:
            fila = _indice_por_paso(_esquina(c)[1], origen_y, paso_y, n_filas)
            celdas.setdefault((fila, col), []).append(c)

    # Orden de lectura: izquierda a derecha, después hacia abajo.
    ordenadas = [celdas.get((f, col), []) for f in range(n_filas) for col in range(n_cols)]
    if not all(ordenadas):
        return None
    # Cada celda tiene que quedarse con exactamente un ancla. Si alguna quedó
    # con dos, la grilla propuesta no es la que tiene el diseño.
    if any(sum(1 for c in g if id(c) in ids_ancla) != 1 for g in ordenadas):
        return None
    return ordenadas


def _detect_slot_bands(components: list[dict]) -> list[list[dict]] | None:
    """Agrupa los componentes de una plantilla multi-producto, un grupo por slot.

    n_slots = GCD de cuántas veces aparece cada variable -- NO el máximo. Una
    variable puede aparecer más de una vez POR slot (un precio partido en
    placeholder de entero + placeholder de decimal, ambos apuntando a la misma
    variable canónica) sin que eso signifique que hay más slots. Con max(), una
    plantilla de 3 slots donde cada precio tiene 2 placeholders se detectaba
    como 6, y de ahí salían grupos mezclando datos de dos productos distintos.

    El conteo mira tanto c["variable"] como c["segments"]: PowerPoint parte los
    placeholders en varios runs al editarlos, y una plantilla real puede tener
    TODOS sus componentes como multi-segmento. Contando solo c["variable"] el
    Counter quedaba vacío y la página entera se trataba como un solo producto.

    La distribución se resuelve como GRILLA (filas x columnas). Los tres
    diseños en uso son distintos:

        3xA4  ->  3 filas x 1 columna   (una debajo de otra)
        A5    ->  1 fila  x 2 columnas  (una al lado de la otra)
        6xA4  ->  3 filas x 2 columnas  (abajo y al costado)

    Ordenando solo por Y --como se hacía antes-- los dos cuadros de una misma
    fila quedan pegados en el orden y el corte los mandaba al mismo grupo: las
    cenefas de la derecha repetían el producto de la izquierda. Se veía en
    6xA4 y en A5; el 3xA4 zafaba por tener una sola columna.

    Devuelve None cuando hay un solo slot (render normal, un producto por
    página).
    """
    import math
    from collections import Counter
    from functools import reduce

    def _comp_variable_names(c: dict) -> list[str]:
        if c.get("variable"):
            return [c["variable"]]
        return [seg["value"] for seg in (c.get("segments") or []) if seg.get("type") == "variable"]

    non_bg = [c for c in components if not c.get("locked")]
    var_counts: Counter = Counter()
    for c in non_bg:
        var_counts.update(_comp_variable_names(c))
    if not var_counts:
        return None

    n_slots = reduce(math.gcd, var_counts.values())
    if n_slots <= 1:
        return None

    # Las ANCLAS son los cuadros de una variable que aparece exactamente una
    # vez por slot (descripcion, codigo...): marcan dónde está cada cenefa.
    ancla = next((v for v, n in var_counts.items() if n == n_slots), None)
    if ancla is not None:
        anclas = [c for c in non_bg if ancla in _comp_variable_names(c)]
        xs_ancla = sorted({round(_esquina(c)[0], 1) for c in anclas})

        # Cuántas columnas hay: se prueba de más a menos, y se acepta la
        # primera grilla que reparta a todos los cuadros dejando un ancla por
        # celda.
        for n_cols in range(min(len(xs_ancla), n_slots), 0, -1):
            if n_slots % n_cols:
                continue
            grilla = _asignar_grilla(non_bg, anclas, n_slots // n_cols, n_cols)
            if grilla is not None:
                return grilla

    # Sin anclas utilizables o con una grilla que no cierra, se cae al criterio
    # viejo: ordenar por Y y cortar en grupos iguales.
    sorted_comps = sorted(non_bg, key=lambda c: _esquina(c)[1])
    total = len(sorted_comps)
    group_size = total // n_slots
    remainder = total % n_slots

    bands: list[list[dict]] = []
    idx = 0
    for i in range(n_slots):
        size = group_size + (1 if i < remainder else 0)
        bands.append(sorted_comps[idx: idx + size])
        idx += size
    return bands


def patch_image_overrides(
    components: list[dict], image_overrides: dict[str, tuple[bytes, str]]
) -> list[dict]:
    """Inyecta imágenes subidas (ej. cocarda) en los componentes de imagen que
    referencien esa variable. Se usa tanto acá como en el paso de preview
    (jobs.py), para que la imagen ya esté horneada en el template_def antes
    de que el usuario llegue a reposicionar."""
    import base64 as _b64

    patched = []
    for c in components:
        var = c.get("variable")
        if c.get("type") == "image" and var and var in image_overrides:
            img_bytes, img_ext = image_overrides[var]
            patched.append({
                **c,
                "image_data": _b64.b64encode(img_bytes).decode(),
                "image_ext":  img_ext,
            })
        else:
            patched.append(c)
    return patched


# ---------------------------------------------------------------------------
# Entry points
# ---------------------------------------------------------------------------

def render_template_to_pptx(
    template_def: dict,
    products: list[dict],
    target_format: str = "a4",
    image_overrides: dict[str, tuple[bytes, str]] | None = None,
    source_pptx_bytes: bytes | None = None,
) -> tuple[bytes, list[str]]:
    """Genera PPTX desde una definición v2 y una lista de productos.

    image_overrides: {variable_name: (image_bytes, ext)} — inyecta imágenes
    subidas en la página de generación como image_data de los componentes
    que usen esa variable.

    source_pptx_bytes: bytes del PPTX original del que se importó el
    template (si vino de un archivo subido, no armado a mano en el editor).
    Cuando está presente, el slide base se reusa/duplica desde ESE archivo
    en vez de reconstruirse desde una presentación en blanco — así se
    preserva el diseño (fondo, master, layout, cualquier shape que el
    importer no haya capturado como componente) tal cual estaba. Sin esto,
    el resultado solo tiene los shapes que SÍ se lograron extraer, y
    cualquier diseño que viva en el layout/master del archivo se pierde.

    Returns (pptx_bytes, missing_vars) donde missing_vars es la lista de
    variables que el template usa pero que no fueron encontradas en el Excel.
    """
    master_format = template_def.get("master_format", "a4")
    components    = template_def.get("components", [])
    rules         = template_def.get("rules", [])

    if image_overrides:
        components = patch_image_overrides(components, image_overrides)
    fmt_info      = get_format(target_format)
    slots         = fmt_info["slots"]

    missing_vars: set[str] = set()

    # ── Detect internal slots from variable repetition count ─────────────────
    # A template encodes N products per slide when its variables each appear N
    # times. Sort components by Y, split into N consecutive groups, and fill
    # each group with one product row — regardless of format or page size.
    slot_bands = _detect_slot_bands(components)

    # Preservar diseño original: solo cuando tenemos los bytes crudos Y el
    # layout no necesita tileo DENTRO de una misma página armado por el
    # motor (slot_bands ya trae sus N celdas pre-armadas en el archivo;
    # slots==1 es una celda por página, tampoco hace falta tilear). El caso
    # restante (una plantilla de una sola celda que el motor debe repetir
    # con offsets dentro de la página, ej. subir un solo pincho y pedirle al
    # sistema que arme la grilla de 6) sigue usando el canvas en blanco.
    preserve_source = source_pptx_bytes is not None and (slot_bands is not None or slots == 1)

    prs = None
    if preserve_source:
        try:
            prs = Presentation(io.BytesIO(source_pptx_bytes))
            if not prs.slides:
                prs = None
        except Exception:
            prs = None
        preserve_source = prs is not None

    if preserve_source:
        # El fondo extraído del MASTER (pptx_importer.py, name="fondo",
        # _source_shape_id=None a propósito) no tiene un shape real en el
        # slide para mutar — cualquier slide que comparta layout/master ya lo
        # hereda visualmente solo con add_slide(), sin dibujar nada de nuevo.
        # Si no se filtra acá, _place_component nunca encuentra shape para
        # mutar y cae al fallback de "crear uno nuevo" en CADA render — y
        # como _duplicate_slide() copia lo que el slide base tenga en ese
        # momento, cada producto siguiente arrastra una copia extra apilada
        # sobre la anterior (bug real, visto con productos duplicando el
        # diseño 2-3 veces encimados).
        components = [
            c for c in components
            if not (c.get("type") == "image" and c.get("variable") is None and c.get("_source_shape_id") is None)
        ]

    if not preserve_source:
        slide_w, slide_h = FORMAT_SLIDES.get(target_format, FORMAT_SLIDES["a4"])
        prs = Presentation()
        prs.slide_width  = slide_w
        prs.slide_height = slide_h

    blank_layout = prs.slide_layouts[6] if not preserve_source else None
    base_slide   = prs.slides[0] if preserve_source else None

    def _next_slide(is_first: bool):
        if preserve_source:
            return base_slide if is_first else _duplicate_slide(prs, base_slide)
        return prs.slides.add_slide(blank_layout)

    if slot_bands:
        bg_comps    = [c for c in components if c.get("locked")]
        n_slots     = len(slot_bands)
        page_groups = [products[i:i + n_slots] for i in range(0, len(products), n_slots)]

        # Todas las hojas se crean ANTES de dibujar. _duplicate_slide copia el
        # slide base tal como está en ese momento, así que clonar después de
        # haber renderizado la página anterior arrastra sus mutaciones -- y
        # ahora que un cuadro vacío se saca del slide, la página 2 nacería sin
        # la cocarda solo porque el primer producto no tenía mecánica.
        hojas = [_next_slide(gi == 0) for gi in range(len(page_groups))]

        for pg, slide in zip(page_groups, hojas):
            shape_map = _shape_id_map(slide.shapes) if preserve_source else {}
            if bg_comps:
                # Los componentes de un slot_bands ya vienen en coordenadas
                # absolutas (una página con N celdas pre-armadas) — nunca hay
                # que re-escalarlos contra target_format, solo posicionarlos
                # tal cual quedaron en master_format. Ver bug histórico: esto
                # antes escalaba por target_format y comprimía la grilla.
                laid_bg = compute_layout(bg_comps, master_format, master_format)
                _render_slide(slide, laid_bg, {}, missing_vars=missing_vars, shape_map=shape_map)
            for band_idx, band_comps in enumerate(slot_bands):
                laid_band = compute_layout(band_comps, master_format, master_format)
                if band_idx < len(pg):
                    product       = pg[band_idx]
                    visibility    = evaluate_rules(rules, product)
                    visible_comps = apply_visibility(laid_band, visibility)
                    visible_comps = _fit_text_to_box(
                        visible_comps, product, get_format(master_format)["width_cm"])
                    _render_slide(slide, visible_comps, product, missing_vars=missing_vars, shape_map=shape_map)
                elif preserve_source:
                    # Página parcial (menos productos que celdas) y estamos
                    # preservando el diseño original: si no se limpia, la
                    # celda sin producto queda con lo que tuviera el archivo
                    # fuente (ej. datos de ejemplo del diseñador).
                    _render_slide(slide, laid_band, {}, shape_map=shape_map)
                # Sin preserve_source: no hay nada que limpiar, la celda
                # simplemente nunca tuvo shapes creados (comportamiento
                # de siempre para el canvas en blanco).

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue(), sorted(missing_vars)

    # ── Single-slot template: one product per format-cell, tiled by offset ───
    # Used for 3xa4 / pinchos / any format where the slide IS the unit cell
    # and multiple cells are arranged spatially on the output page. Cuando
    # preserve_source es True acá slots siempre es 1 (gateado más arriba),
    # así que el offset por slot da siempre 0 — un producto, un slide, sin
    # tileo interno.
    laid_out  = compute_layout(components, target_format, master_format)
    slot_cols = fmt_info.get("slot_cols", 1)
    cell_w    = fmt_info["width_cm"]
    cell_h    = fmt_info["height_cm"]
    groups    = [products[i:i + slots] for i in range(0, len(products), slots)]

    for gi, group in enumerate(groups):
        slide     = _next_slide(gi == 0)
        shape_map = _shape_id_map(slide.shapes) if preserve_source else {}

        for slot_idx, product in enumerate(group):
            col = slot_idx % slot_cols
            row = slot_idx // slot_cols
            slot_offset_x = col * cell_w
            slot_offset_y = row * cell_h

            visibility    = evaluate_rules(rules, product)
            visible_comps = apply_visibility(laid_out, visibility)
            # El offset del slot corre el cuadro dentro de la hoja, asi que el
            # ancho de papel disponible se mide desde donde va a caer de verdad.
            visible_comps = _fit_text_to_box(
                visible_comps, product, max(0.5, fmt_info["width_cm"] * slot_cols - slot_offset_x))

            _render_slide(slide, visible_comps, product, slot_offset_x, slot_offset_y, missing_vars=missing_vars, shape_map=shape_map)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), sorted(missing_vars)


def generate_from_template_v2(
    template_def: dict,
    excel_bytes: bytes,
    target_format: str = "a4",
    vigencia: str = "",
    legales: str = "",
    usar_legales: bool = False,
    image_overrides: dict[str, tuple[bytes, str]] | None = None,
) -> tuple[bytes, list[str]]:
    """Parsea Excel y genera PPTX desde una definición de template.

    Camino directo sin jobs -- lo usa la generación sincrónica. El flujo
    normal de la plataforma pasa por jobs.py (preview + confirmación).
    """
    products = load_products_from_bytes(excel_bytes, vigencia, legales, usar_legales)
    return render_template_to_pptx(template_def, products, target_format, image_overrides)
