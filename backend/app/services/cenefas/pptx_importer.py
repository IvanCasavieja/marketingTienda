"""Importador PPTX → definición v2 de componentes."""
import base64
import re
import uuid
from io import BytesIO

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

_EMU_PER_CM = 360_000

_RE_PLACEHOLDER = re.compile(r"<<(\w+?)(\d*)>>", re.IGNORECASE)

# ---------------------------------------------------------------------------
# Mapa de placeholders → (variable_name, tipo, transform)
#
# variable_name: nombre canónico camelCase almacenado en el JSON del template
# tipo:          "text" | "price" | "image"
# transform:     "smart_bold" | "upper" | "price_full" | "none"
#
# Los placeholders en el PPTX deben ser <<variableName>> (mismos nombres que
# las columnas del Excel). Se soportan nombres legacy por backward compat.
# ---------------------------------------------------------------------------
_PLACEHOLDER_MAP: dict[str, tuple[str, str, str]] = {
    # ── Nombres canónicos nuevos ──────────────────────────────────────────
    "precioactual":       ("precioActual",      "price", "price_full"),
    "precioanterior":     ("precioAnterior",    "price", "price_full"),
    "preciobanco":        ("precioBanco",       "price", "price_full"),
    "banco":              ("banco",             "text",  "none"),
    "descripcion":        ("descripcion",       "text",  "smart_bold"),
    "titulo":             ("mecanica",           "text",  "upper"),
    "aclaracion":         ("aclaracion",        "text",  "none"),
    "aclaracion1":        ("aclaracion",        "text",  "none"),  # alias — misma variable que "aclaracion"
    "aclaracion2":        ("aclaracion2",       "text",  "none"),
    "aclaracion3":        ("aclaracion3",       "text",  "none"),
    "segundaaclaracion":  ("segundaAclaracion", "text",  "none"),
    "vigencia":           ("vigencia",          "text",  "none"),
    "codigosku":          ("codigoSKU",         "text",  "none"),
    "dia":                ("dia",               "text",  "upper"),
    "mes":                ("mes",               "text",  "none"),
    "ano":                ("año",               "text",  "none"),
    "año":                ("año",               "text",  "none"),
    "moneda":             ("moneda",            "text",  "none"),
    "categoria":          ("categoria",         "text",  "none"),
    "subcategoria":       ("subCategoria",      "text",  "none"),
    "descuento":          ("descuento",         "text",  "none"),

    # ── Legacy backward compat (plantillas importadas antes de la unificación) ──
    "p":               ("precioActual",      "price", "price_full"),
    "precio":          ("precioActual",      "price", "price_full"),
    "pbanco":          ("precioBanco",       "price", "price_full"),
    "p1":              ("mecanica",           "text",  "upper"),
    "mecanica":        ("mecanica",          "text",  "upper"),
    "code":            ("codigoSKU",         "text",  "none"),
    "otraaclaracion":  ("segundaAclaracion", "text",  "none"),
    "unidadprecio":    ("unidadPrecio",      "text",  "none"),
    "unidadpbanco":    ("unidadPBanco",      "text",  "none"),
    "unidad":          ("unidadPrecio",      "text",  "none"),
    "unidadmedida":    ("unidadPrecio",      "text",  "none"),
    "descripci":       ("descripcion",       "text",  "smart_bold"),

    # ── Parrilla y Vinos: precio principal + 3 niveles por cantidad ───────
    # (4x3/5x3/6x3), cada nivel partido en placeholder de entero + placeholder
    # de decimal -- ver _ALIASES_PARRILLA_VINOS_OVERRIDE en data_engine.py
    # para el lado del Excel. precioP (no precioActual) a propósito: el "$"
    # de esta plantilla ya es texto fijo en el PPTX, precioActual vendría
    # con su propio "$" auto-formateado y quedaría duplicado ("$$399").
    "preciop":         ("precioP",           "price", "price_full"),
    "4x3p":            ("precio4x3",         "price", "price_integer"),
    "decimal4x3":      ("precio4x3",         "price", "price_decimal"),
    "5x3p":            ("precio5x3",         "price", "price_integer"),
    "decimal5x3":      ("precio5x3",         "price", "price_decimal"),
    "6x3p":            ("precio6x3",         "price", "price_integer"),
    "decimal6x3":      ("precio6x3",         "price", "price_decimal"),
}

# Mapa: variable_name → columna Excel canónica (para el panel de Variables)
_CSV_COLUMN_MAP: dict[str, str] = {
    "precioActual":      "precioActual",
    "precioAnterior":    "precioAnterior",
    "precioBanco":       "precioBanco",
    "banco":             "banco",
    "descripcion":       "descripcion",
    "mecanica":          "mecanica",
    "aclaracion":        "aclaracion",
    "aclaracion2":       "aclaracion2",
    "aclaracion3":       "aclaracion3",
    "segundaAclaracion": "segundaAclaracion",
    "vigencia":          "vigencia",
    "codigoSKU":         "codigoSKU",
    "dia":               "dia",
    "mes":               "mes",
    "año":               "año",
    "moneda":            "moneda",
    "categoria":         "categoria",
    "subCategoria":      "subCategoria",
    "descuento":         "descuento",
    # legacy
    "mecanica":          "mecanica",
    "unidadPrecio":      "unidadPrecio",
    "unidadPBanco":      "unidadPBanco",
}

_REQUIRED_VARS = {"descripcion", "precioActual"}

_FORMATS_DIM = {
    "a4":      (21.0,  29.7,  1),
    "a3":      (29.7,  42.0,  1),
    "3xa4":    (21.0,  9.9,   3),   # slide de una franja; 3 en A4 portrait
    "pinchos": (7.0,   14.85, 6),   # slide de un pincho; grilla 3×2 en A4
    "a5":      (14.85, 21.0,  1),
    "6xa4":    (7.0,   14.85, 6),   # slide de una celda; grilla 3×2 en A4 (arte propio)
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _emu_to_cm(emu: int) -> float:
    return round(emu / _EMU_PER_CM, 3)


def _detect_format(width_cm: float, height_cm: float) -> tuple[str, int]:
    best = "a4"
    best_dist = float("inf")
    for fmt_id, (w, h, _slots) in _FORMATS_DIM.items():
        dist = abs(width_cm - w) + abs(height_cm - h)
        if dist < best_dist:
            best_dist = dist
            best = fmt_id
    _, _, slots = _FORMATS_DIM[best]
    return best, slots


_MAX_IMAGE_BYTES = 300_000  # ~300 KB antes de comprimir

# Formatos que los navegadores entienden nativamente
_WEB_EXTS = {"jpeg", "jpg", "png", "gif", "webp"}


def _extract_image_b64(shape) -> tuple[str, str] | None:
    """Extrae la imagen de un shape Picture como base64 web-compatible.
    WMF/EMF/BMP/TIFF se convierten a PNG/JPEG con Pillow.
    Devuelve (base64_str, ext) o None."""
    try:
        img_obj = shape.image
        raw = img_obj.blob
        ext = img_obj.ext.lower().lstrip(".")
        if ext == "jpg":
            ext = "jpeg"

        if ext not in _WEB_EXTS:
            # Intentar convertir a PNG (funciona en Windows/GDI+)
            converted = _to_web_image(raw)
            if converted is not None:
                raw, ext = converted
            # Si PIL falla (Linux): guardar raw con ext original.
            # El renderer lo embebe en el PPTX directamente sin PIL.
        elif len(raw) > _MAX_IMAGE_BYTES:
            raw, ext = _compress_image(raw, ext)

        return base64.b64encode(raw).decode("utf-8"), ext
    except Exception:
        return None


def _to_web_image(raw: bytes) -> tuple[bytes, str] | None:
    """Intenta convertir WMF/EMF/BMP/TIFF → PNG usando PIL.
    Funciona en Windows (GDI+). En Linux devuelve None → el renderer embebe raw."""
    try:
        from PIL import Image as PILImage
        import io as _io

        img = PILImage.open(_io.BytesIO(raw))

        # Verificar que PIL realmente pudo rasterizar (no solo abrir el header)
        img.load()

        max_dim = 800
        if max(img.width, img.height) > max_dim:
            ratio = max_dim / max(img.width, img.height)
            img = img.resize(
                (int(img.width * ratio), int(img.height * ratio)),
                PILImage.LANCZOS,
            )

        has_alpha = img.mode in ("RGBA", "LA") or (
            img.mode == "P" and "transparency" in img.info
        )
        buf = _io.BytesIO()
        img.convert("RGBA" if has_alpha else "RGB").save(buf, format="PNG", optimize=True)
        return buf.getvalue(), "png"
    except Exception:
        return None


def _compress_image(raw: bytes, ext: str) -> tuple[bytes, str]:
    """Comprime una imagen web existente (JPEG/PNG). Preserva transparencia."""
    try:
        from PIL import Image as PILImage
        import io as _io

        img = PILImage.open(_io.BytesIO(raw))
        has_alpha = img.mode in ("RGBA", "LA") or (
            img.mode == "P" and "transparency" in img.info
        )

        max_dim = 1500
        if max(img.width, img.height) > max_dim:
            ratio = max_dim / max(img.width, img.height)
            img = img.resize(
                (int(img.width * ratio), int(img.height * ratio)),
                PILImage.LANCZOS,
            )

        buf = _io.BytesIO()
        if has_alpha:
            img.convert("RGBA").save(buf, format="PNG", optimize=True)
            return buf.getvalue(), "png"
        else:
            img.convert("RGB").save(buf, format="JPEG", quality=85, optimize=True)
            return buf.getvalue(), "jpeg"
    except Exception:
        return raw, ext


def _extract_fill_color(shape) -> str | None:
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        # Try python-pptx direct RGB (works for explicit fills)
        try:
            rgb = fill.fore_color.rgb
            return f"#{rgb.red:02X}{rgb.green:02X}{rgb.blue:02X}"
        except (TypeError, AttributeError):
            pass
        # Fallback: parse XML for srgbClr (explicit RGB embedded in XML)
        from lxml import etree
        from pptx.oxml.ns import qn
        solidFill = shape.element.find('.//' + qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                val = srgbClr.get('val', '')
                if val:
                    return f"#{val.upper()}"
            # Theme color: try to resolve from presentation theme XML
            schemeClr = solidFill.find(qn('a:schemeClr'))
            if schemeClr is not None:
                return _resolve_theme_color(shape, schemeClr)
        return None
    except Exception:
        return None


def _resolve_theme_color(shape, schemeClr) -> str | None:
    """Resolve a theme color reference (schemeClr) to an explicit hex RGB."""
    try:
        from pptx.oxml.ns import qn
        val = schemeClr.get('val', '')
        _COLOR_ELEM = {
            'dk1': 'a:dk1', 'lt1': 'a:lt1',
            'dk2': 'a:dk2', 'lt2': 'a:lt2',
            'accent1': 'a:accent1', 'accent2': 'a:accent2',
            'accent3': 'a:accent3', 'accent4': 'a:accent4',
            'accent5': 'a:accent5', 'accent6': 'a:accent6',
        }
        tag = _COLOR_ELEM.get(val)
        if not tag:
            return None
        # Walk up to find the theme part
        part = shape.part
        theme_part = None
        for rel in part.rels.values():
            if 'theme' in rel.reltype:
                theme_part = rel._target
                break
        if theme_part is None:
            # Try slide layout → slide master → theme
            try:
                layout = part.slide_layout
                for rel in layout.slide_master.part.rels.values():
                    if 'theme' in rel.reltype:
                        theme_part = rel._target
                        break
            except Exception:
                pass
        if theme_part is None:
            return None
        theme_el = theme_part._element
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        # Path: a:theme/a:themeElements/a:clrScheme/a:accent1/a:srgbClr
        node = theme_el.find(f'.//a:clrScheme/{tag}/a:srgbClr', ns)
        if node is None:
            node = theme_el.find(f'.//a:clrScheme/{tag}/a:sysClr', ns)
        if node is not None:
            hex_val = node.get('val') or node.get('lastClr')
            if hex_val:
                # Apply luminance modifier if present
                lum_mod = schemeClr.find(qn('a:lumMod'))
                lum_off = schemeClr.find(qn('a:lumOff'))
                r = int(hex_val[0:2], 16)
                g = int(hex_val[2:4], 16)
                b = int(hex_val[4:6], 16)
                if lum_mod is not None:
                    factor = int(lum_mod.get('val', '100000')) / 100000
                    r = int(r * factor); g = int(g * factor); b = int(b * factor)
                if lum_off is not None:
                    offset = int(lum_off.get('val', '0')) / 100000 * 255
                    r = int(min(255, r + offset))
                    g = int(min(255, g + offset))
                    b = int(min(255, b + offset))
                return f"#{r:02X}{g:02X}{b:02X}"
    except Exception:
        pass
    return None


def _resolve_scheme_colors(slide) -> dict[str, str]:
    """Resuelve clrScheme + clrMap del theme del master a hex reales — sin
    esto, cualquier run con color de tema (ej. texto blanco a propósito sobre
    un fondo oscuro, algo muy común en diseños reales) se pierde silenciosamente
    al recrearse: python-pptx expone el color de tema como un slot ("bg1",
    "accent1"...), nunca como RGB directo, así que hay que ir a buscar el RGB
    al theme.xml del master nosotros mismos."""
    try:
        from pptx.oxml.ns import qn as _qn

        master = slide.slide_layout.slide_master
        theme_part = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        from lxml import etree
        theme_el = etree.fromstring(theme_part.blob)
        clr_scheme_el = theme_el.find(".//" + _qn("a:clrScheme"))
        scheme: dict[str, str] = {}
        if clr_scheme_el is not None:
            for child in clr_scheme_el:
                tag = child.tag.split("}")[-1]
                srgb = child.find(_qn("a:srgbClr"))
                sysclr = child.find(_qn("a:sysClr"))
                val = srgb.get("val") if srgb is not None else (sysclr.get("lastClr") if sysclr is not None else None)
                if val:
                    scheme[tag] = val

        # bg1/tx1/bg2/tx2 pasan por un mapa de indirección propio del master
        # (clrMap) antes de resolver contra el clrScheme — accentN/hlink no.
        clr_map_el = master._element.find(_qn("p:clrMap"))
        clr_map = dict(clr_map_el.attrib) if clr_map_el is not None else {}

        resolved: dict[str, str] = dict(scheme)
        for slot, target in clr_map.items():
            if target in scheme:
                resolved[slot] = scheme[target]
        return resolved
    except Exception:
        return {}


def _extract_font_color(run, theme_colors: dict[str, str] | None = None) -> str | None:
    try:
        from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR

        color = run.font.color
        if color.type == MSO_COLOR_TYPE.SCHEME:
            if theme_colors:
                slot = MSO_THEME_COLOR.to_xml(color.theme_color)
                hexval = theme_colors.get(slot)
                if hexval:
                    return f"#{hexval}"
            return None
        # RGBColor es una subclase de tuple con __str__ propio que ya da el
        # hex de 6 dígitos ("2E323E") -- NO tiene atributos .red/.green/.blue
        # (esos no existen en python-pptx, ver pptx.dml.color.RGBColor).
        # El acceso a esos atributos tiraba AttributeError silenciosamente
        # tragado por el except de abajo -- por eso NINGÚN color explícito
        # se extraía nunca, en ningún destino, desde que existe esta función.
        return f"#{color.rgb}"
    except Exception:
        return None


def _get_shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)


def _resolve_placeholder(root: str, suffix: str) -> tuple[str, str, str]:
    """Resuelve un match de _RE_PLACEHOLDER (root en minúscula + sufijo
    numérico opcional) a (variable_name, tipo, transform) — compartido entre
    el caso de un solo placeholder por shape y el de varios (_build_segments)."""
    # Reconstruct full name including numeric suffix (e.g. "aclaracion" + "2" → "aclaracion2")
    full = root + suffix
    if full in _PLACEHOLDER_MAP:
        return _PLACEHOLDER_MAP[full]
    if root in _PLACEHOLDER_MAP:
        return _PLACEHOLDER_MAP[root]
    for key, value in _PLACEHOLDER_MAP.items():
        if full.startswith(key) or key.startswith(full):
            return value
    return (full, "text", "none")


def _detect_placeholder(text: str) -> tuple[str, str, str] | None:
    m = _RE_PLACEHOLDER.search(text)
    if not m:
        return None
    return _resolve_placeholder(m.group(1).lower(), m.group(2))


def _run_style(run, theme_colors: dict[str, str] | None = None) -> dict:
    """Estilo (tamaño, negrita, tachado, familia, color) de UN run puntual —
    usado por _build_segments para que cada segmento de un cuadro con texto
    mixto conserve su propio formato, en vez de heredar uniformemente el
    estilo "de shape" que arma _extract_style a partir del primer run."""
    style: dict = {}
    font = run.font
    try:
        if font.size:
            style["font_size"] = round(font.size.pt, 1)
    except Exception:
        pass
    try:
        b = font.bold
        if b is None:
            from pptx.oxml.ns import qn as _qn_b
            rPr = run._r.find(_qn_b("a:rPr"))
            if rPr is not None:
                b_str = rPr.get("b")
                if b_str is not None:
                    b = b_str not in ("0", "false")
        if b is not None:
            style["font_bold"] = bool(b)
    except Exception:
        pass
    try:
        from pptx.oxml.ns import qn as _qn_s
        rPr = run._r.find(_qn_s("a:rPr"))
        if rPr is not None:
            strike = rPr.get("strike")
            if strike is not None and strike != "noStrike":
                style["strikethrough"] = True
    except Exception:
        pass
    try:
        if font.name:
            style["font_family"] = font.name
    except Exception:
        pass
    color = _extract_font_color(run, theme_colors)
    if color:
        style["color"] = color
    return style


def _build_segments(
    shape, text: str, matches: list, theme_colors: dict[str, str] | None = None
) -> tuple[list[dict], list[tuple[str, str]]]:
    """Parte el texto de un shape con texto estático + variable mezclados en
    segmentos que se resuelven de forma independiente, cada uno con el
    estilo real del run que lo originó — ej. un diseñador escribió
    "<<Aclaracion1>> <<Aclaracion2>> <<Aclaracion3>>" todo en un solo cuadro
    en vez de 3 cuadros separados (_detect_placeholder con un solo match se
    quedaba con el primero y perdía el resto), o "$" + "<<Precio>>" con
    tamaños de fuente distintos en la misma caja de precio de Rompe Precios
    (antes se perdía el tamaño real del número, todo el cuadro quedaba con
    el tamaño del primer run). El mismo <<...>> que delimita cada variable
    en el resto del sistema es lo que se usa acá para encontrar los límites
    de cada segmento; el estilo de cada uno sale del run real que lo cubre,
    no de un único estilo "de shape"."""
    # Mapea cada run a su rango [start, end) en el texto aplanado, en el
    # mismo orden de concatenación que usa _get_shape_text() — pero medido
    # sobre el texto SIN recortar (raw), porque _parse_shape le hace un
    # .strip() antes de llegar acá y eso desplaza las posiciones.
    raw = _get_shape_text(shape)
    left_trim = len(raw) - len(raw.lstrip())

    spans = []
    pos = 0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            start = pos - left_trim
            pos += len(run.text)
            spans.append((start, pos - left_trim, run))

    def _style_at(start: int, end: int) -> dict:
        for s, e, run in spans:
            if s < end and e > start:
                return _run_style(run, theme_colors)
        return {}

    segments: list[dict] = []
    seg_vars: list[tuple[str, str]] = []
    pos = 0
    for m in matches:
        if m.start() > pos:
            segments.append({"type": "static", "value": text[pos:m.start()], "style": _style_at(pos, m.start())})
        var_name, var_type, transform = _resolve_placeholder(m.group(1).lower(), m.group(2))
        segments.append({"type": "variable", "value": var_name, "transform": transform, "style": _style_at(m.start(), m.end())})
        seg_vars.append((var_name, var_type))
        pos = m.end()
    if pos < len(text):
        segments.append({"type": "static", "value": text[pos:], "style": _style_at(pos, len(text))})
    return segments, seg_vars


def _extract_style(shape, theme_colors: dict[str, str] | None = None) -> dict:
    style: dict = {}
    if not shape.has_text_frame:
        return style
    tf = shape.text_frame
    for para in tf.paragraphs:
        if para.alignment is not None:
            _align = {PP_ALIGN.LEFT: "left", PP_ALIGN.CENTER: "center", PP_ALIGN.RIGHT: "right"}
            style["align"] = _align.get(para.alignment, "center")
            break
    first_run = None
    max_font_size = 0.0
    for para in tf.paragraphs:
        for run in para.runs:
            try:
                if run.font.size:
                    max_font_size = max(max_font_size, run.font.size.pt)
            except Exception:
                pass
            if first_run is None and run.text.strip():
                first_run = run
    if first_run:
        font = first_run.font
        try:
            if font.size:
                style["font_size"] = round(font.size.pt, 1)
        except Exception:
            pass
        try:
            b = font.bold
            if b is None:
                # font.bold is None when bold is inherited from theme/layout.
                # Fall back to reading the XML directly.
                from pptx.oxml.ns import qn as _qn_b
                rPr = first_run._r.find(_qn_b("a:rPr"))
                if rPr is not None:
                    b_str = rPr.get("b")
                    if b_str is not None:
                        b = b_str not in ("0", "false")
                if b is None:
                    # Check paragraph-level defRPr as last resort
                    for _para in tf.paragraphs:
                        if any(r._r is first_run._r for r in _para.runs):
                            def_rpr = _para._p.find(_qn_b("a:defRPr"))
                            if def_rpr is not None:
                                b_str = def_rpr.get("b")
                                if b_str is not None:
                                    b = b_str not in ("0", "false")
                            break
            if b is not None:
                style["font_bold"] = bool(b)
        except Exception:
            pass
        try:
            from pptx.oxml.ns import qn as _qn_s
            rPr = first_run._r.find(_qn_s("a:rPr"))
            if rPr is not None:
                strike = rPr.get("strike")
                if strike is not None and strike != "noStrike":
                    style["strikethrough"] = True
        except Exception:
            pass
        try:
            if font.name:
                style["font_family"] = font.name
        except Exception:
            pass
        color = _extract_font_color(first_run, theme_colors)
        if color:
            style["color"] = color
    # Max font across all runs (spacer runs set line height)
    if max_font_size > style.get("font_size", 0):
        style["line_height_pt"] = round(max_font_size, 1)
    style.setdefault("align", "center")

    # Vertical anchor + autofit + insets from bodyPr
    try:
        from pptx.oxml.ns import qn as _qn
        body_pr = tf._txBody.find(_qn("a:bodyPr"))
        if body_pr is not None:
            anchor = body_pr.get("anchor", "t")
            if anchor and anchor != "t":
                style["vertical_align"] = anchor
            has_norm = body_pr.find(_qn("a:normAutofit")) is not None
            has_sp   = body_pr.find(_qn("a:spAutoFit"))  is not None
            style["auto_fit"] = has_norm or has_sp
            # Bottom inset (needed for anchor=b overflow calculation)
            style["_b_ins_emu"] = int(body_pr.get("bIns", "45720"))
        else:
            style["auto_fit"] = False
    except Exception:
        style["auto_fit"] = False

    # Line spacing from lstStyle (e.g. 90000 = 90%)
    try:
        from pptx.oxml.ns import qn as _qn2
        lst = tf._txBody.find(_qn2("a:lstStyle"))
        if lst is not None:
            pct_el = lst.find(f'.//{_qn2("a:lnSpc")}/{_qn2("a:spcPct")}')
            if pct_el is not None:
                style["_lnSpc_pct"] = int(pct_el.get("val", "100000")) / 1000
    except Exception:
        pass

    # Baseline shift from first non-empty run (e.g. 30000 = 30% superscript)
    try:
        from pptx.oxml.ns import qn as _qn3
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    rPr = run._r.find(_qn3("a:rPr"))
                    if rPr is not None:
                        bl = rPr.get("baseline")
                        if bl and int(bl) != 0:
                            style["_baseline"] = int(bl)
                    break
            if "_baseline" in style:
                break
    except Exception:
        pass

    return style


def _flatten_shapes(shapes) -> list:
    result = []
    for shape in shapes:
        if hasattr(shape, "shapes"):
            result.extend(_flatten_shapes(shape.shapes))
        else:
            result.append(shape)
    return result


def _make_common(shape, z_index: int) -> dict | None:
    try:
        left   = _emu_to_cm(shape.left   or 0)
        top    = _emu_to_cm(shape.top    or 0)
        width  = _emu_to_cm(shape.width  or 0)
        height = _emu_to_cm(shape.height or 0)
    except Exception:
        return None
    if width < 0.1 or height < 0.1:
        return None
    return {
        "id":               str(uuid.uuid4()),
        # id del shape original en el PPTX fuente — permite que el render
        # final mute ESE shape en vez de reconstruir todo desde cero, así se
        # preserva el master/layout/diseño del archivo (ver component_renderer).
        "_source_shape_id": shape.shape_id,
        "base_bounds":      {"x": left, "y": top, "width": width, "height": height},
        "format_overrides": {},
        "z_index":          z_index,
        "locked":           False,
        "visible":          True,
    }


# ---------------------------------------------------------------------------
# Parseo de shapes individuales
# ---------------------------------------------------------------------------

def _parse_shape(
    shape, z_index: int, theme_colors: dict[str, str] | None = None,
    allow_single_placeholder_segments: bool = False,
) -> dict | None:
    common = _make_common(shape, z_index)
    if common is None:
        return None

    # Imagen embebida (foto, cocarde, logo)
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            comp = {**common, "type": "image", "name": f"imagen_{z_index}",
                    "variable": None, "style": {}, "_var_type": "image_url"}
            result = _extract_image_b64(shape)
            if result:
                comp["image_data"], comp["image_ext"] = result
            return comp
    except Exception:
        pass

    # Shape con texto
    if shape.has_text_frame:
        text = _get_shape_text(shape).strip()
        style = _extract_style(shape, theme_colors)

        # Recompute bounds for bottom-anchored shapes with spacer overflow runs.
        # PowerPoint positions these using anchor=b + large spacer + baseline shift.
        # We convert to a top-anchored box at the computed visual position instead,
        # so the export faithfully places the text where it appears in the original.
        if (style.get("vertical_align") == "b"
                and style.get("line_height_pt", 0) > common["base_bounds"]["height"] * 1.5):
            max_font   = style["line_height_pt"]
            lnSpc      = style.pop("_lnSpc_pct", 100) / 100
            b_ins      = style.pop("_b_ins_emu", 45720) / _EMU_PER_CM
            baseline   = style.pop("_baseline", 0) / 100000   # e.g. 0.30
            eff_line_h = max_font * lnSpc / 72 * 2.54          # cm
            y_bottom   = common["base_bounds"]["y"] + common["base_bounds"]["height"] - b_ins
            bl_shift   = baseline * max_font / 72 * 2.54       # cm upward
            new_y      = y_bottom - eff_line_h - bl_shift
            common["base_bounds"]["y"]      = round(new_y, 3)
            common["base_bounds"]["height"] = round(eff_line_h, 3)
            style.pop("vertical_align", None)   # use default top anchor
            style.pop("line_height_pt", None)   # no spacer run needed in export
        else:
            style.pop("_lnSpc_pct", None)
            style.pop("_b_ins_emu", None)
            style.pop("_baseline", None)

        placeholder_matches = list(_RE_PLACEHOLDER.finditer(text))
        non_empty_runs = [r for p in shape.text_frame.paragraphs for r in p.runs if r.text.strip()]
        # 2+ placeholders en un cuadro (aclaracion1/2/3 compartiendo caja,
        # comportamiento preexistente, válido en cualquier plantilla) O un
        # solo placeholder pero con texto estático de otro run al lado (ej.
        # "$" + "<<Precio>>" con tamaños distintos) — esto último es nuevo y
        # se restringe a destinos que lo pidieron explícitamente
        # (allow_single_placeholder_segments, hoy solo Rompe Precios) para no
        # cambiar el resultado de ninguna plantilla existente sin que se
        # vuelva a importar a propósito.
        multi_run_trigger = allow_single_placeholder_segments and len(non_empty_runs) > 1
        if placeholder_matches and (len(placeholder_matches) > 1 or multi_run_trigger):
            segments, seg_vars = _build_segments(shape, text, placeholder_matches, theme_colors)
            return {**common, "type": "text", "name": (text[:30] or "texto"),
                    "variable": None, "segments": segments,
                    "style": style, "_seg_vars": seg_vars}

        ph = _detect_placeholder(text)
        if ph:
            var_name, var_type, transform = ph
            # TODO(cocarda/imagen): un <<imagen>> escrito como texto (no una
            # imagen insertada de verdad en el shape) SIEMPRE cae acá con
            # type="text" — más abajo, _parse_shape solo produce type="image"
            # para shapes cuyo shape_type ya es MSO_SHAPE_TYPE.PICTURE. Como
            # patch_image_overrides() (component_renderer.py) solo parchea
            # componentes con type=="image", la cocarda de Rompe Precios NUNCA
            # se aplica mientras el pptx tenga el placeholder como texto — hay
            # que pedirle al diseñador que inserte una imagen real (Insertar >
            # Imágenes) en esa posición. Pendiente: decidir si vale la pena
            # detectar var_name=="imagen" acá y forzar type="image" con un
            # placeholder gris hasta que se reemplace, en vez de solo avisar
            # por missing_vars como hace hoy.
            return {**common, "type": "text", "name": var_name,
                    "variable": var_name, "transform": transform,
                    "style": style, "_var_type": var_type}
        if text:
            label = text[:30]
            return {**common, "type": "text", "name": label,
                    "variable": None, "static_value": text,
                    "transform": "none", "style": style, "_var_type": "text"}
        # Texto vacío: puede ser un shape decorativo con color de fondo → fall-through

    # Shape sin texto (o texto vacío) → fondo/decorativo con color
    style = {}
    bg = _extract_fill_color(shape)
    if bg:
        style["background_color"] = bg
        return {**common, "type": "shape", "name": f"fondo_{z_index}", "style": style}
    return None


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def import_pptx(pptx_bytes: bytes, name: str = "Template importado", category: str | None = None) -> dict:
    """Parsea el primer slide de un PPTX y devuelve una definición v2.

    category: destino declarado por el caller (ej. "rompe_precios") — hoy
    solo se usa para habilitar la detección de segmentos con un único
    placeholder + texto estático (ver _parse_shape), restringida a Rompe
    Precios y Parrilla y Vinos (mismo patrón de plantilla, mismo pedido
    explícito de comportarse igual) a pedido explícito para no cambiar el
    import de ninguna otra plantilla existente."""
    allow_single_placeholder_segments = category in ("rompe_precios", "parrilla_y_vinos")
    prs = Presentation(BytesIO(pptx_bytes))
    if not prs.slides:
        raise ValueError("El archivo PPTX no tiene slides")

    slide = prs.slides[0]
    width_cm  = _emu_to_cm(prs.slide_width)
    height_cm = _emu_to_cm(prs.slide_height)
    format_id, slots = _detect_format(width_cm, height_cm)
    slot_width = width_cm / slots

    components: list[dict]          = []
    variables_seen: dict[str, dict] = {}
    z_index = 0
    theme_colors = _resolve_scheme_colors(slide)

    # ── 1. Imágenes del slide master (fondo visual del template) ──────────
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        master = slide.slide_layout.slide_master
        for shape in master.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            common = _make_common(shape, z_index)
            if common is None:
                continue
            result = _extract_image_b64(shape)
            if not result:
                continue
            b64, ext = result
            components.append({
                **common,
                "type":       "image",
                "name":       "fondo",
                "variable":   None,
                # None a propósito (pisa lo que puso _make_common): este
                # shape_id es del MASTER, un namespace de ids totalmente
                # aparte del slide — nunca va a matchear un shape real ahí, y
                # si coincidiera por casualidad con el id de otro shape del
                # slide sería peor (mutaría el shape equivocado). El render
                # (component_renderer.py) usa este None como señal explícita
                # de "fondo heredado del master, no hay nada que dibujar de
                # nuevo cuando se preserva el diseño original".
                "_source_shape_id": None,
                "image_data": b64,
                "image_ext":  ext,
                "style":      {},
                "locked":     True,
            })
            z_index += 1
    except Exception:
        pass

    # ── 2. Shapes del slide (datos variables + imágenes embebidas) ────────
    for shape in _flatten_shapes(slide.shapes):
        comp = _parse_shape(shape, z_index, theme_colors, allow_single_placeholder_segments)
        if comp is None:
            continue

        # Formatos multi-slot horizontales (pinchos): solo primera columna
        if slots > 1 and comp["base_bounds"]["x"] >= slot_width:
            continue

        components.append(comp)
        z_index += 1

        var_name = comp.get("variable")
        if var_name and var_name not in variables_seen:
            variables_seen[var_name] = {
                "type":       comp.pop("_var_type", "text"),
                "csv_column": _CSV_COLUMN_MAP.get(var_name, var_name.upper()),
            }
        else:
            comp.pop("_var_type", None)

        # Componente multi-segmento (varios <<...>> en un mismo cuadro) — cada
        # variable que usa también tiene que quedar en el panel de Variables.
        for seg_var, seg_type in comp.pop("_seg_vars", []):
            if seg_var not in variables_seen:
                variables_seen[seg_var] = {
                    "type":       seg_type,
                    "csv_column": _CSV_COLUMN_MAP.get(seg_var, seg_var.upper()),
                }

    variables = [
        {
            "name":       vname,
            "type":       vinfo["type"],
            "required":   vname in _REQUIRED_VARS,
            "csv_column": vinfo["csv_column"],
        }
        for vname, vinfo in variables_seen.items()
    ]

    return {
        "version":       "2.0",
        "name":          name,
        "master_format": format_id,
        "formats":       [format_id],
        "variables":     variables,
        "components":    components,
        "rules":         [],
    }
