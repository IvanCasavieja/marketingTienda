"""Motor de renderizado PPTX — genera presentaciones desde datos de productos."""
import copy
import io
import math
import re

from lxml import etree
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

from app.services.cenefas.data_engine import load_products_from_bytes
from app.services.cenefas.formatters import (
    split_caps,
    P1_FONT_SIZE,
    P1_BOLD,
    P1_MARGIN_EMU,
    PRICE_SYMBOL_PT,
    PRICE_INT_PT,
    PRICE_DECIMAL_PT,
    PBANCO_INT_PT,
    UNIDAD_PRECIO_PT,
    UNIDAD_PBANCO_PT,
    DESC_PT,
    DESC_MIN_PT,
)

# ---------------------------------------------------------------------------
# Helpers de texto
# ---------------------------------------------------------------------------

def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)


def _set_text(shape, text: str) -> None:
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        runs = para.runs
        if runs:
            runs[0].text = text
            for run in runs[1:]:
                run.text = ""
            return
    if shape.text_frame.paragraphs:
        para = shape.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = text


def _set_text_sized(shape, text: str, pt: int) -> None:
    _set_text(shape, text)
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(pt)


def _set_price(shape, text: str, int_pt: int | None = None, preserve_sizes: bool = False) -> None:
    """Renderiza precio con símbolo, entero y decimal en runs separados.

    preserve_sizes=True → no toca atributos sz; respeta los tamaños del template original.
    """
    if not shape.has_text_frame:
        return

    if text.startswith("U$S "):
        symbol, number = "U$S ", text[4:]
    elif text.startswith("$"):
        symbol, number = "$", text[1:]
    else:
        _set_text(shape, text)
        return

    target_para = None
    for para in shape.text_frame.paragraphs:
        if para.runs:
            target_para = para
            break
    if target_para is None:
        return

    p_elem = target_para._p
    tmpl_r_sym = copy.deepcopy(target_para.runs[0]._r)
    tmpl_r_int = copy.deepcopy(target_para.runs[-1]._r)

    for r_elem in list(p_elem.findall(qn("a:r"))):
        p_elem.remove(r_elem)

    if "," in number:
        num_int, num_dec_digits = number.rsplit(",", 1)
        num_dec = "," + num_dec_digits
    else:
        num_int = number
        num_dec = None

    # Calcular tamaños de fuente ANTES de construir los runs.
    # Si int_pt is None (modo A4), achica símbolo + entero + decimal proporcionalmente
    # para que el ancho total quede dentro de shape.width.
    _sym_pt_final = PRICE_SYMBOL_PT
    _int_pt_final = int_pt if int_pt is not None else PRICE_INT_PT
    _dec_pt_final = PRICE_DECIMAL_PT

    if not preserve_sizes and int_pt is None and shape.width > 0:
        K = 12700 * 0.58
        dec_chars = len(num_dec) if num_dec else 0
        total_w = (
            len(symbol)  * PRICE_SYMBOL_PT  +
            len(num_int) * PRICE_INT_PT      +
            dec_chars    * PRICE_DECIMAL_PT
        ) * K
        if total_w > shape.width and total_w > 0:
            scale = shape.width / total_w
            _sym_pt_final = max(int(PRICE_SYMBOL_PT  * scale), 8)
            _int_pt_final = max(int(PRICE_INT_PT     * scale), PBANCO_INT_PT)
            _dec_pt_final = max(int(PRICE_DECIMAL_PT * scale), 8)

    sym_r = copy.deepcopy(tmpl_r_sym)
    sym_r.find(qn("a:t")).text = symbol
    if not preserve_sizes and int_pt is None:
        _rpr = sym_r.find(qn("a:rPr"))
        if _rpr is not None:
            _rpr.set("sz", str(_sym_pt_final * 100))

    int_r = copy.deepcopy(tmpl_r_int)
    int_r.find(qn("a:t")).text = num_int
    if not preserve_sizes:
        _rpr = int_r.find(qn("a:rPr"))
        if _rpr is not None:
            _rpr.set("sz", str(_int_pt_final * 100))

    runs = [sym_r, int_r]

    if num_dec:
        dec_r = copy.deepcopy(tmpl_r_int)
        dec_r.find(qn("a:t")).text = num_dec
        if not preserve_sizes:
            dec_rPr = dec_r.find(qn("a:rPr"))
            if dec_rPr is not None:
                dec_rPr.set("sz", str(_dec_pt_final * 100))
        runs.append(dec_r)

    end_rpr = p_elem.find(qn("a:endParaRPr"))
    if end_rpr is not None:
        idx = list(p_elem).index(end_rpr)
        for r in reversed(runs):
            p_elem.insert(idx, r)
    else:
        for r in runs:
            p_elem.append(r)


def _set_desc(shape, text: str, preserve_sizes: bool = False) -> None:
    if not shape.has_text_frame:
        return
    if not text:
        _set_text(shape, "")
        return
    tf = shape.text_frame
    target_para = None
    for para in tf.paragraphs:
        if para.runs:
            target_para = para
            break
    if target_para is None:
        return

    parts = split_caps(text)
    first_seg, first_bold = parts[0]
    first_run = target_para.runs[0]
    tmpl_r = copy.deepcopy(first_run._r)
    if not preserve_sizes:
        # A4 Redex: forzar tamaño DESC_PT
        _tmpl_rPr = tmpl_r.find(qn("a:rPr"))
        if _tmpl_rPr is not None:
            _tmpl_rPr.set("sz", str(DESC_PT * 100))
        first_run.font.size = Pt(DESC_PT)
    first_run.text = first_seg
    first_run.font.bold = first_bold

    p_elem = target_para._p
    all_r = list(p_elem.findall(qn("a:r")))
    for r_elem in all_r[1:]:
        p_elem.remove(r_elem)

    for seg, bold in parts[1:]:
        if not seg:
            continue
        new_r = copy.deepcopy(tmpl_r)
        new_r.find(qn("a:t")).text = seg
        rPr = new_r.find(qn("a:rPr"))
        if rPr is not None:
            rPr.set("b", "1" if bold else "0")
        end_rpr = p_elem.find(qn("a:endParaRPr"))
        if end_rpr is not None:
            p_elem.insert(list(p_elem).index(end_rpr), new_r)
        else:
            p_elem.append(new_r)


def _set_normAutofit(shape) -> None:
    """Cambia spAutoFit → normAutofit para que el texto se achique en vez de expandir el cuadro."""
    if not shape.has_text_frame:
        return
    body_pr = shape.text_frame._txBody.find(qn("a:bodyPr"))
    if body_pr is None:
        return
    for tag in (qn("a:spAutoFit"), qn("a:noAutofit")):
        child = body_pr.find(tag)
        if child is not None:
            body_pr.remove(child)
    if body_pr.find(qn("a:normAutofit")) is None:
        body_pr.append(etree.Element(qn("a:normAutofit")))


def _set_p1(shape, text: str) -> None:
    _set_text(shape, text)
    if not shape.has_text_frame:
        return
    is_combo_label = bool(re.match(r"^\d+X$", text.strip()))
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if not is_combo_label:
                run.font.size = Pt(P1_FONT_SIZE)
            run.font.bold = P1_BOLD


def _is_multi_sku(code: str) -> bool:
    return bool(code and ("/" in code or re.search(r"\d\s*[-–—]\s*\d", code)))


def _estimate_text_height_emu(text: str, shape_width_emu: int, font_pt: int) -> int:
    """Estimación aproximada de altura de texto en EMU."""
    if not text or shape_width_emu <= 0 or font_pt <= 0:
        return int(font_pt * 12700 * 1.4)
    char_w = max(1, int(font_pt * 12700 * 0.5))
    lines  = math.ceil(len(text) / max(1, shape_width_emu // char_w))
    return lines * int(font_pt * 12700 * 1.4)


def _set_runs_font_size(shape, pt: int) -> None:
    """Cambia el tamaño de fuente de todos los runs de un shape."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(pt)


# ---------------------------------------------------------------------------
# Llenado de slots
# ---------------------------------------------------------------------------

def _apply_desc_lateral_margin(shape, slide_width: int, margin_emu: int = 360000) -> None:
    """Restringe el shape de descripción para que no se salga de los márgenes laterales del slide.

    Solo se aplica en modo A4 (1 producto/slide). margin_emu default = 1 cm.
    """
    if slide_width <= 0:
        return
    new_left  = max(shape.left, margin_emu)
    new_right = min(shape.left + shape.width, slide_width - margin_emu)
    if new_right > new_left:
        shape.left  = new_left
        shape.width = new_right - new_left


def _fill_slot(shapes, data: dict, adjust_p1: bool = True, slide_height: int = 0, slide_width: int = 0, a4_mode: bool = True) -> None:
    """Rellena un slot de cenefa con los datos del producto.

    a4_mode=True  → A4 Redex: sobreescribe tamaños de letra y aplica ajustes de layout.
    a4_mode=False → Pass-through: respeta exactamente el template original (3xA4, Pinchos, etc.).
    """
    expanded = list(shapes)
    for shape in list(shapes):
        if hasattr(shape, "shapes"):
            expanded.extend(shape.shapes)

    p1_shape         = None
    price_shape      = None
    oferta_shape     = None
    desc_shape       = None
    aclaracion_shape = None
    code  = data.get("codigoSKU", "")
    multi = _is_multi_sku(code)
    ps    = not a4_mode  # preserve_sizes — True cuando NO es A4

    for shape in expanded:
        t = _shape_text(shape)
        if re.search(r"<<P\d+>>", t):
            p1_shape = shape
            _set_p1(shape, data.get("mecanica", ""))
        elif re.search(r"Precio\s+\d+", t) or re.search(r"<<Precio\d*>>", t):
            price_shape = shape
            if a4_mode:
                _apply_desc_lateral_margin(shape, slide_width)
            _set_price(shape, data.get("precioActual", ""), preserve_sizes=ps)
            if a4_mode:
                _set_normAutofit(shape)
        elif re.search(r"<<Mecanica\d+>>", t):
            _set_text(shape, data.get("mecanica", ""))
        elif "<<" in t and "Descripci" in t:
            desc_shape = shape
            _set_desc(shape, data.get("descripcion", ""), preserve_sizes=ps)
            if a4_mode:
                _set_normAutofit(shape)
                _apply_desc_lateral_margin(shape, slide_width)
        elif re.search(r"<<UnidadMedida\d+>>", t):
            _set_text(shape, "unidad" if multi else "")
        elif re.search(r"<<Vigencia\d*>>", t):
            _set_text(shape, data.get("vigencia", ""))
        elif re.search(r"<<Aclaracion\d*>>", t):
            aclaracion_shape = shape
            _set_text(shape, data.get("aclaracion", ""))
        elif re.search(r"<<OtraAclaracion\d*>>", t):
            _set_text(shape, data.get("segundaAclaracion", ""))
        elif re.search(r"<<[Dd][Ii][Aa]\d*>>", t):
            _set_text(shape, data.get("dia", ""))
        elif re.search(r"<<oferta\d*>>", t, re.IGNORECASE):
            oferta_shape = shape
            _set_text(shape, data.get("oferta", ""))
        elif re.search(r"<<[Cc]ode\d*>>", t):
            _set_text(shape, code)
        elif re.search(r"<<[Pp][Bb]anco\d*>>", t, re.IGNORECASE):
            _set_price(shape, data.get("precioBanco", ""), int_pt=PBANCO_INT_PT, preserve_sizes=ps)
        elif re.search(r"<<[Bb]anco\d*>>", t, re.IGNORECASE):
            _set_text(shape, data.get("banco", ""))
        elif re.search(r"<<UnidadPrecio\d*>>", t):
            _set_text_sized(shape, "unidad" if multi else "", UNIDAD_PRECIO_PT)
        elif re.search(r"<<UnidadPBanco\d*>>", t):
            _set_text_sized(shape, "unidad" if multi else "", UNIDAD_PBANCO_PT)
        elif t.strip().lower() == "unidad":
            _set_text(shape, "unidad" if multi else "")

    # ═══════════════════════════════════════════════════════════════════════════
    # A4 REDEX — ajustes de layout exclusivos para cenefa de 1 producto/slide
    # ═══════════════════════════════════════════════════════════════════════════
    if a4_mode:
        # MxN también tiene oferta="" pero tiene mecanica; solo precio fijo tiene ambas vacías
        is_precio_fijo = (
            not data.get("oferta", "").strip() and
            not data.get("mecanica", "").strip()
        )

        if is_precio_fijo and oferta_shape is not None and price_shape is not None:
            # Precio fijo: sube precio y descripción para ocupar el espacio vacío de oferta
            shift = oferta_shape.height
            price_shape.top -= shift
            if desc_shape is not None:
                desc_shape.top -= shift

        if adjust_p1 and p1_shape is not None and price_shape is not None:
            p1_shape.top = price_shape.top - P1_MARGIN_EMU

        # Combo/MxN: detecta colisión descripción vs bases y condiciones
        if not is_precio_fijo and desc_shape is not None and aclaracion_shape is not None:
            desc_text = data.get("descripcion", "")
            est_h = _estimate_text_height_emu(desc_text, desc_shape.width, DESC_PT)
            if desc_shape.top + est_h > aclaracion_shape.top:
                overlap = (desc_shape.top + est_h) - aclaracion_shape.top
                new_top = aclaracion_shape.top + overlap
                if slide_height > 0 and new_top + aclaracion_shape.height <= slide_height:
                    aclaracion_shape.top = new_top
                else:
                    font_pt = DESC_PT - 2
                    while font_pt >= DESC_MIN_PT:
                        if desc_shape.top + _estimate_text_height_emu(desc_text, desc_shape.width, font_pt) <= aclaracion_shape.top:
                            break
                        font_pt -= 2
                    _set_runs_font_size(desc_shape, font_pt)
    # ═══════════════════════════════════════════════════════════════════════════


def _clear_slot(shapes) -> None:
    for shape in shapes:
        if not shape.has_text_frame:
            continue
        t = _shape_text(shape)
        if "<<" in t or re.search(r"Precio\s+\d", t):
            _set_text(shape, "")


# ---------------------------------------------------------------------------
# Detección de slots
# ---------------------------------------------------------------------------

def _slot_num(shape) -> int | None:
    t = _shape_text(shape)
    for n in range(1, 10):
        if (f"<<P{n}>>" in t
                or f"Precio {n}" in t
                or f"<<Precio{n}>>" in t
                or f"<<Mecanica{n}>>" in t
                or f"<<UnidadMedida{n}>>" in t
                or f"<<Vigencia{n}>>" in t
                or f"<<Aclaracion{n}>>" in t
                or f"<<OtraAclaracion{n}>>" in t
                or f"<<UnidadPrecio{n}>>" in t
                or f"<<UnidadPBanco{n}>>" in t
                or bool(re.search(rf"<<[Dd]escripci[oó]n{n}>>", t))
                or bool(re.search(rf"<<[Dd][Ii][Aa]{n}>>", t))):
            return n
    return None


def _buckets_coherent(buckets: dict, present: list) -> bool:
    """True si las shapes de cada bucket están en la misma región visual (spread vertical < 3M EMU)."""
    for n in present:
        shapes_in_bucket = buckets[n]
        if len(shapes_in_bucket) < 2:
            continue
        tops = [s.top for s in shapes_in_bucket]
        if max(tops) - min(tops) > 3_000_000:
            return False
    return True


def _assign_by_nearest_anchor(anchors: list, all_shapes: list) -> list[list]:
    """Asigna cada shape al slot cuyo anchor es espacialmente más cercano (para plantillas Pinchos)."""
    products = len(anchors)
    anchors_sorted = sorted(
        anchors,
        key=lambda s: (s.top + s.height // 2, s.left + s.width // 2),
    )

    def _ctr(s):
        return (s.left + s.width / 2, s.top + s.height / 2)

    anchor_centers = [_ctr(a) for a in anchors_sorted]
    slots: list[list] = [[] for _ in range(products)]
    for shape in all_shapes:
        cx, cy = _ctr(shape)
        best = min(
            range(products),
            key=lambda i: (cx - anchor_centers[i][0]) ** 2 + (cy - anchor_centers[i][1]) ** 2,
        )
        slots[best].append(shape)
    return slots


def _get_slots(shapes) -> list[list]:
    """Agrupa shapes por slot de producto. Soporta 1–9 productos por slide."""
    buckets: dict[int, list] = {n: [] for n in range(1, 10)}
    for shape in shapes:
        n = _slot_num(shape)
        if n:
            buckets[n].append(shape)

    present = sorted(n for n in range(1, 10) if buckets[n])

    if len(present) > 1:
        if _buckets_coherent(buckets, present):
            return [buckets[n] for n in present]
        all_shapes = list(shapes)
        anchors = [s for s in all_shapes if re.search(r"Precio\s+\d+|<<Precio\d*>>", _shape_text(s))]
        if not anchors:
            anchors = [s for s in all_shapes if re.search(r"<<P\d+>>", _shape_text(s))]
        if anchors:
            return _assign_by_nearest_anchor(anchors, all_shapes)
        return [buckets[n] for n in present]

    all_shapes = list(shapes)
    anchors = [s for s in all_shapes if re.search(r"<<P\d+>>", _shape_text(s))]
    if not anchors:
        anchors = [s for s in all_shapes if re.search(r"Precio\s+\d+|<<Precio\d*>>", _shape_text(s))]

    products = len(anchors)
    if products <= 1:
        return [all_shapes]

    xs = [s.left for s in anchors]
    ys = [s.top  for s in anchors]
    horizontal = (max(xs) - min(xs)) >= (max(ys) - min(ys))

    key = (lambda s: s.left + s.width // 2) if horizontal else (lambda s: s.top + s.height // 2)
    sorted_shapes = sorted(all_shapes, key=key)
    per_slot  = len(sorted_shapes) // products
    slots     = [sorted_shapes[i * per_slot:(i + 1) * per_slot] for i in range(products)]
    remainder = sorted_shapes[products * per_slot:]
    if remainder:
        slots[-1].extend(remainder)
    return slots


# ---------------------------------------------------------------------------
# Layout A4
# ---------------------------------------------------------------------------

def _center_content_a4(slide, slide_width: int) -> None:
    """Centra shapes de texto en slides A4 sin GroupShapes."""
    if any(hasattr(s, "shapes") for s in slide.shapes):
        return
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        shape.left  = 0
        shape.width = slide_width
        for para in shape.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER


def _align_bank_group_a4(slide) -> None:
    """Recorta el ancho del precio para que no invada la columna del GROUP bancario."""
    price_shape = None
    group_shape = None
    for shape in slide.shapes:
        if hasattr(shape, "shapes"):
            group_shape = shape
        elif shape.has_text_frame:
            t = _shape_text(shape)
            if re.search(r"<<Precio\d*>>", t) or re.search(r"Precio\s+\d+", t):
                price_shape = shape
    if price_shape is None or group_shape is None:
        return
    gap = 150000  # ~4mm
    price_right_max = group_shape.left - gap
    if price_shape.left + price_shape.width > price_right_max:
        price_shape.width = price_right_max - price_shape.left


# ---------------------------------------------------------------------------
# Construcción de slides
# ---------------------------------------------------------------------------

def _duplicate_slide(prs, template_slide):
    """Duplica un slide completo: shapes, fondo e imágenes con sus relaciones.

    Modifica el spTree en-place en vez de reemplazar el cSld completo para evitar
    invalidar el lazyproperty cache de slide.shapes que python-pptx construye
    durante add_slide. El cache sigue apuntando al mismo nodo spTree, que ahora
    contiene los children del template.
    """
    R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    new_slide = prs.slides.add_slide(template_slide.slide_layout)

    # Copiar todas las relaciones del template al nuevo slide (excepto slideLayout)
    rId_map = {}
    for rId, rel in template_slide.part.rels.items():
        if rel.reltype.endswith("/slideLayout"):
            continue
        try:
            new_rId = new_slide.part.relate_to(
                rel._target, rel.reltype, is_external=rel.is_external
            )
            rId_map[rId] = new_rId
        except Exception:
            pass

    def _remap_rids(elem):
        for el in elem.iter():
            for attr_key in list(el.attrib):
                ns_part, _, local = attr_key.partition("}")
                ns_part = ns_part.lstrip("{")
                if ns_part == R_NS and local in ("embed", "link", "id"):
                    old_rId = el.attrib[attr_key]
                    if old_rId in rId_map:
                        el.attrib[attr_key] = rId_map[old_rId]
        return elem

    template_cSld = template_slide._element.find(qn("p:cSld"))
    if template_cSld is None:
        return new_slide

    new_cSld = new_slide._element.find(qn("p:cSld"))
    if new_cSld is None:
        return new_slide

    # Copiar fondo (p:bg) si existe en el template
    template_bg = template_cSld.find(qn("p:bg"))
    if template_bg is not None:
        existing_bg = new_cSld.find(qn("p:bg"))
        if existing_bg is not None:
            new_cSld.remove(existing_bg)
        new_spTree_ref = new_cSld.find(qn("p:spTree"))
        bg_pos = list(new_cSld).index(new_spTree_ref) if new_spTree_ref is not None else 0
        new_cSld.insert(bg_pos, _remap_rids(copy.deepcopy(template_bg)))

    # Reemplazar children del spTree en-place: limpia los placeholders vacíos
    # que add_slide copió desde el layout y agrega los shapes del template
    template_spTree = template_cSld.find(qn("p:spTree"))
    new_spTree = new_cSld.find(qn("p:spTree"))

    if template_spTree is not None and new_spTree is not None:
        for child in list(new_spTree):
            new_spTree.remove(child)
        for child in list(template_spTree):
            new_spTree.append(_remap_rids(copy.deepcopy(child)))

    return new_slide


# ---------------------------------------------------------------------------
# Margen horizontal
# ---------------------------------------------------------------------------

def _apply_horizontal_margin(slide, slide_width: int, margin_emu: int) -> None:
    """Contrae todos los shapes para que quepan dentro del margen izquierdo/derecho."""
    if margin_emu <= 0:
        return
    for shape in slide.shapes:
        original_left  = shape.left
        original_right = shape.left + shape.width
        new_left  = max(original_left,  margin_emu)
        new_right = min(original_right, slide_width - margin_emu)
        if new_right > new_left:
            shape.left  = new_left
            shape.width = new_right - new_left


# ---------------------------------------------------------------------------
# Entry point principal
# ---------------------------------------------------------------------------

def generate_pptx_bytes(
    excel_bytes: bytes,
    template_bytes: bytes,
    vigencia: str,
    aclaracion: str,
    otra_alcohol: str,
    banco: str = "",
    margin_cm: float = 0.0,
) -> bytes:
    products = load_products_from_bytes(excel_bytes, vigencia, aclaracion, otra_alcohol, banco)

    prs = Presentation(io.BytesIO(template_bytes))
    if not prs.slides:
        raise ValueError("La plantilla PPTX está vacía.")

    template_slide     = prs.slides[0] if len(prs.slides) == 1 else prs.slides[1]
    initial_slots      = _get_slots(list(template_slide.shapes))
    products_per_slide = max(len(initial_slots), 1)
    slide_height       = prs.slide_height
    margin_emu         = int(margin_cm * 914400 / 2.54) if margin_cm else 0

    groups = [products[i:i + products_per_slide] for i in range(0, len(products), products_per_slide)]

    # Pre-crear todos los slides como duplicados del template ANTES de llenar ninguno.
    # Esto es crítico: si llenáramos template_slide primero (idx=0), las copias
    # posteriores tendrían datos del producto 1 en vez de los placeholders originales.
    all_slides = [template_slide]
    for _ in range(len(groups) - 1):
        all_slides.append(_duplicate_slide(prs, template_slide))

    for slide, group in zip(all_slides, groups):
        cur_slots = _get_slots(list(slide.shapes))
        is_a4 = products_per_slide == 1
        for i, product in enumerate(group):
            if i < len(cur_slots):
                _fill_slot(cur_slots[i], product, adjust_p1=is_a4, slide_height=slide_height, slide_width=prs.slide_width, a4_mode=is_a4)
        for i in range(len(group), products_per_slide):
            if i < len(cur_slots):
                _clear_slot(cur_slots[i])
        if products_per_slide == 1:
            _center_content_a4(slide, prs.slide_width)
            _align_bank_group_a4(slide)
        if margin_emu > 0:
            _apply_horizontal_margin(slide, prs.slide_width, margin_emu)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
