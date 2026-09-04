"""Fija el importer de PPTX y el render — los bugs visuales de esta semana.

Los PPTX de prueba se arman en memoria con python-pptx: no hay archivos
binarios en el repo ni dependencia de las plantillas reales de la base.
"""
import io
import re
import zipfile

from pptx import Presentation
from pptx.util import Cm, Pt

from app.services.cenefas.component_renderer import (
    _ancho_disponible_cm,
    render_template_to_pptx,
)
from app.services.cenefas.pptx_importer import import_pptx


def _pptx_con_textos(*textos, size_pt=100):
    """Un A4 con un cuadro de texto por cada string recibido."""
    prs = Presentation()
    prs.slide_width = Cm(21.0)
    prs.slide_height = Cm(29.7)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i, texto in enumerate(textos):
        box = slide.shapes.add_textbox(Cm(2), Cm(3 + i * 5), Cm(17), Cm(4))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = texto
        run.font.size = Pt(size_pt)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _componente(defin, contiene):
    return next(c for c in defin["components"] if contiene in str(c))


# ---------------------------------------------------------------------------
# Importer
# ---------------------------------------------------------------------------

def test_texto_fijo_junto_al_placeholder_no_se_pierde():
    # Bug real (28/08): "$<<precioOferta>>" en UN solo run descartaba el "$"
    # en silencio — la A4 REDEX re-subida con el símbolo salía igual que antes.
    d = import_pptx(_pptx_con_textos("$<<precioOferta>>"))
    comp = _componente(d, "precioOferta")
    segs = comp.get("segments") or []
    assert any(s["type"] == "static" and "$" in s["value"] for s in segs)
    assert any(s["type"] == "variable" and s["value"] == "precioOferta" for s in segs)


def test_placeholder_canonico_directo():
    d = import_pptx(_pptx_con_textos("<<precioOferta>>"))
    comp = _componente(d, "precioOferta")
    assert comp.get("variable") == "precioOferta"
    assert not comp.get("segments")


def test_moneda_legacy_resuelve_unidad_moneda():
    # Desde el 29/08 el símbolo es variable de nuevo: un <<Moneda>> viejo
    # apunta a unidadMoneda en vez de importarse como cuadro vacío.
    d = import_pptx(_pptx_con_textos("<<Moneda>>"))
    assert any(c.get("variable") == "unidadMoneda" for c in d["components"])


def test_placeholder_dado_de_baja_queda_vacio():
    d = import_pptx(_pptx_con_textos("<<UnidadMedida1>>"))
    comp = d["components"][0]
    assert comp.get("variable") is None
    assert comp.get("static_value") == ""


def test_dos_placeholders_en_un_cuadro():
    d = import_pptx(_pptx_con_textos("<<unidadMoneda>><<precioOferta>>"))
    comp = _componente(d, "precioOferta")
    vars_ = [s["value"] for s in comp["segments"] if s["type"] == "variable"]
    assert vars_ == ["unidadMoneda", "precioOferta"]


# ---------------------------------------------------------------------------
# Render de punta a punta (sin base, sin red)
# ---------------------------------------------------------------------------

def _runs_del_pptx(pptx_bytes):
    with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as z:
        xml = "".join(
            z.read(n).decode("utf-8", "ignore")
            for n in z.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", n))
    return "".join(re.findall(r"<a:t>([^<]*)</a:t>", xml))


def test_render_imprime_simbolo_y_precio():
    src = _pptx_con_textos("$<<precioOferta>>", "<<descripcion>>")
    d = import_pptx(src)
    pptx, missing = render_template_to_pptx(
        d, [{"precioOferta": "1.290", "descripcion": "Aspiradora MASTER-X"}],
        "a4", None, src)
    runs = _runs_del_pptx(pptx)
    assert "$" in runs and "1.290" in runs and "MASTER-X" in runs


def test_render_dolares():
    src = _pptx_con_textos("<<unidadMoneda>><<precioOferta>>")
    d = import_pptx(src)
    pptx, _ = render_template_to_pptx(
        d, [{"unidadMoneda": "U$S", "precioOferta": "60"}], "a4", None, src)
    assert "U$S" in _runs_del_pptx(pptx)


def test_variable_ausente_en_el_excel_se_reporta():
    # El cuadro tiene texto fijo ("$") ademas de la variable: se dibuja igual
    # y la columna ausente se reporta. (Un cuadro SOLO de variables vacias se
    # saca entero del slide antes de llegar a reportar nada -- es la regla
    # que evita imprimir la cocarda de color sin contenido.)
    src = _pptx_con_textos("$<<precioOferta>>")
    d = import_pptx(src)
    _, missing = render_template_to_pptx(d, [{"descripcion": "X"}], "a4", None, src)
    assert "precioOferta" in missing


# ---------------------------------------------------------------------------
# El achique por vecinos — el solapamiento del 3xA4 (29/08)
# ---------------------------------------------------------------------------

def _caja(x, y, w, h, variable, texto=""):
    return {
        "id": variable, "type": "text", "variable": variable,
        "style": {"font_size": 97.0},
        "base_bounds": {"x": x, "y": y, "width": w, "height": h},
        "computed_bounds": {"x": x, "y": y, "width": w, "height": h},
    }


def test_una_vecina_chata_tambien_limita_el_ancho():
    # Bug real (29/08): la caja del precio del 3xA4 mide 8,55 de alto y la de
    # la descripción 1,08 — la descripción la pisa ENTERA pero era menos que
    # la mitad de 8,55, no contaba como vecina, y "U$S449" se imprimía encima
    # del texto. El umbral es la mitad del alto DEL MÁS BAJO de los dos.
    precio = _caja(1.68, 4.53, 19.83, 8.55, "precioOferta")
    desc = _caja(9.38, 4.86, 8.42, 1.08, "descripcion")
    producto = {"precioOferta": "449", "descripcion": "Refrigerador MIDEA"}
    disponible = _ancho_disponible_cm(precio, [precio, desc], producto, 21.0)
    # limitado por la descripción (arranca en 9,38), no por su propia caja
    assert disponible < 9.38 - 1.68
    # y una vecina VACÍA no limita nada
    disponible_sola = _ancho_disponible_cm(
        precio, [precio, desc], {"precioOferta": "449", "descripcion": ""}, 21.0)
    assert disponible_sola > 9.38 - 1.68


def test_etiqueta_comprando_no_limita_el_precio_que_decora():
    # Caso real (Preciazos-202608-A4, producción): "Comprando 2"
    # (tipoOfertaComprando) vive superpuesta arriba de precioOferta a
    # propósito -- geométricamente es CASI IDÉNTICA a la "vecina chata" de
    # arriba (46,7% del ancho vs 42,5%, 99,5% de solape vertical vs 100%):
    # ninguna relación de ancho/alto/solape sola alcanza para distinguirlas.
    # La señal que sí distingue es la variable: "descripcion" es contenido de
    # verdad, "tipoOfertaComprando" es una etiqueta pensada para flotar sobre
    # un precio (ver _VARIABLES_ETIQUETA_FLOTANTE). Sin ese chequeo por
    # variable, este caso rompía el fix de la "vecina chata": cualquier
    # heurística puramente geométrica que lo arreglara volvía a tratar a
    # "Comprando 2" como pared y achicaba el precio al mínimo.
    precio = _caja(3.559, 7.589, 17.5, 7.951, "precioOferta")
    comprando = _caja(6.262, 7.574, 8.177, 2.821, "tipoOfertaComprando")
    producto = {"precioOferta": "129", "tipoOfertaComprando": "Comprando 2"}
    disponible = _ancho_disponible_cm(precio, [precio, comprando], producto, 21.0)
    assert disponible > 6.262 - 3.559


def test_mxn_imprime_el_literal_una_sola_vez():
    # Bug real (pag. 54 de mundo hogar): la A4 REDEX tiene cocarda
    # (tipoOferta) Y cuadro que tapa al precio (promoOferta) -- en un M x N
    # los dos llevan el mismo literal y "2X1" salia impreso dos veces. La
    # regla de excluyentes esconde la cocarda cuando promoOferta tapa.
    src = _pptx_con_textos("<<tipoOferta>>", "<<unidadMoneda>><<precioOferta>>", "<<promoOferta>>")
    d = import_pptx(src)
    pptx, _ = render_template_to_pptx(
        d, [{"tipoOferta": "2x1", "promoOferta": "2x1",
             "unidadMoneda": "$", "precioOferta": "49,50"}],
        "a4", None, src)
    runs = _runs_del_pptx(pptx)
    assert runs.count("2x1") == 1, f"el literal salio {runs.count('2x1')} veces: {runs!r}"
    # y el precio quedo tapado: no se imprime
    assert "49,50" not in runs


def test_combo_muestra_el_precio_con_cocarda():
    # En un combo promoOferta viene vacia -> el precio unitario se ve, con la
    # cocarda arriba, aunque el diseno tenga el cuadro de promoOferta.
    src = _pptx_con_textos("<<tipoOferta>>", "<<unidadMoneda>><<precioOferta>>", "<<promoOferta>>")
    d = import_pptx(src)
    pptx, _ = render_template_to_pptx(
        d, [{"tipoOferta": "2x$299", "promoOferta": "",
             "unidadMoneda": "$", "precioOferta": "149,50"}],
        "a4", None, src)
    runs = _runs_del_pptx(pptx)
    assert "2x$299" in runs and "149,50" in runs and "$" in runs
